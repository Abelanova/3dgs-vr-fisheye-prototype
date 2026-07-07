from pathlib import Path
from zipfile import ZipFile, ZIP_DEFLATED
from tempfile import NamedTemporaryFile
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from PIL import Image

import matplotlib
matplotlib.use("Agg")
from matplotlib import rcParams
import matplotlib.pyplot as plt

rcParams["mathtext.fontset"] = "stix"
rcParams["font.family"] = "STIXGeneral"


ROOT = Path(r"E:\3dgs fisheye projection")
SRC = Path(r"C:\Users\nova\Downloads\FOV and Fisheye Distortion Rendering of 3D Gaussian Splats for VR Scene Exploration (2).pptx")
OUT = ROOT / "Presentation" / "Revised_FOV_Fisheye_and_Splat_Footprint_2slides_v2.pptx"
ASSETS = ROOT / "Presentation" / "revised_two_slide_assets"


BLACK = RGBColor(16, 16, 16)
GREY = RGBColor(92, 92, 92)
LIGHT_GREY = RGBColor(247, 248, 250)
MID = RGBColor(210, 218, 228)
TEAL = RGBColor(0, 135, 145)
BLUE = RGBColor(43, 93, 173)
RED = RGBColor(196, 65, 62)
PURPLE = RGBColor(113, 82, 165)
GOLD = RGBColor(186, 142, 45)
GREEN = RGBColor(64, 145, 98)


def render_formula(expr: str, out: Path, fontsize=25, color="#202020"):
    out.parent.mkdir(parents=True, exist_ok=True)
    fig = plt.figure(figsize=(0.01, 0.01), dpi=320)
    text = fig.text(0, 0, expr, fontsize=fontsize, color=color)
    fig.canvas.draw()
    bbox = text.get_window_extent()
    w = max(bbox.width / fig.dpi + 0.18, 0.35)
    h = max(bbox.height / fig.dpi + 0.16, 0.25)
    plt.close(fig)

    fig = plt.figure(figsize=(w, h), dpi=320)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")
    ax.text(0.5, 0.5, expr, ha="center", va="center", fontsize=fontsize, color=color)
    fig.savefig(out, transparent=True, bbox_inches="tight", pad_inches=0.02)
    plt.close(fig)


def fit_size(image_path, max_w, max_h=None):
    with Image.open(image_path) as im:
        px_w, px_h = im.size
    aspect = px_h / px_w
    fit_w = max_w
    fit_h = fit_w * aspect
    if max_h is not None and fit_h > max_h:
        fit_h = max_h
        fit_w = fit_h / aspect
    return fit_w, fit_h


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_wrapped_paragraph(slide, x, y, w, h, lines, size=12.5, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = line
        r.font.name = "Aptos"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, 0.72, 0.48, 11.7, 0.48, title, 27, True)
    if subtitle:
        add_text(slide, 0.74, 0.98, 11.2, 0.28, subtitle, 11, False, GREY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.28), Inches(0.7), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def add_footer(slide, num):
    add_text(slide, 0.72, 7.08, 6.4, 0.18, "3DGS FOV/Fisheye Rendering for VR Scene Exploration", 6.8, False, RGBColor(130, 130, 130))
    add_text(slide, 12.2, 7.08, 0.4, 0.18, str(num), 7, False, RGBColor(130, 130, 130), PP_ALIGN.RIGHT)


def add_formula(slide, image_path, x, y, w, h=None, caption=None):
    fit_w, fit_h = fit_size(image_path, w, h)
    x_offset = (w - fit_w) / 2 if w > fit_w else 0
    pic = slide.shapes.add_picture(str(image_path), Inches(x + x_offset), Inches(y), width=Inches(fit_w), height=Inches(fit_h))
    try:
        pic.shadow.inherit = False
    except Exception:
        pass
    if caption:
        add_text(slide, x, y + fit_h + 0.03, w, 0.16, caption, 7.2, False, GREY, PP_ALIGN.CENTER)
    return pic


def rect_label(slide, x, y, w, h, text, line_color, fill=RGBColor(255, 255, 255), size=11):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line_color
    sh.line.width = Pt(1.1)
    tf = sh.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = BLACK
    return sh


def arrow(slide, x1, y1, x2, y2, color=GREY):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.3)
    conn.line.end_arrowhead = True
    return conn


def strip_visual_effects_from_pptx(path: Path):
    effect_tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)
    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in effect_tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(path)


def extract_slide8_images():
    ASSETS.mkdir(parents=True, exist_ok=True)
    src_prs = Presentation(SRC)
    slide = src_prs.slides[7]
    imgs = []
    for idx, shape in enumerate(slide.shapes):
        if hasattr(shape, "image"):
            img = shape.image
            ext = img.ext
            path = ASSETS / f"slide8_image_{idx}.{ext}"
            path.write_bytes(img.blob)
            imgs.append((shape.left, path))
    imgs.sort(key=lambda t: t[0])
    # The first four picture objects are the web comparison images in the source slide.
    return [p for _, p in imgs[:4]]


def build():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    ASSETS.mkdir(parents=True, exist_ok=True)

    formulas = {
        "fov": (r"$f_y=\frac{1}{\tan(\mathrm{FOV}_y/2)}$", 27),
        "theta": (r"$\theta=\operatorname{atan2}\!\left(\sqrt{x^2+y^2},-z\right),\quad r=g(\theta)$", 24),
        "cov": (r"$\Sigma_{2D}=J\,W\,\Sigma_{3D}\,W^{T}J^{T}$", 28),
        "fish_cov": (r"$\Sigma_{2D}^{fish}=J_{fish}\,W\,\Sigma_{3D}\,W^{T}J_{fish}^{T}$", 27),
    }
    rendered = {}
    for k, (expr, size) in formulas.items():
        path = ASSETS / f"{k}.png"
        render_formula(expr, path, size)
        rendered[k] = path

    imgs = extract_slide8_images()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1: FOV vs fisheye.
    s = prs.slides.add_slide(blank)
    add_title(s, "Wide FOV vs Nonlinear Fisheye Projection",
              "Rectilinear FOV expands the camera frustum; fisheye maps ray angle to screen radius.")
    x_positions = [0.72, 3.42, 7.05, 9.75]
    labels = ["FOV 180°", "Fisheye 180°", "FOV 360°", "Fisheye 360°"]
    for x, img, lab in zip(x_positions, imgs, labels):
        add_text(s, x, 1.55, 2.15, 0.18, lab, 9, True, BLACK)
        fit_w, fit_h = fit_size(img, 2.25, 1.33)
        s.shapes.add_picture(str(img), Inches(x), Inches(1.77), width=Inches(fit_w), height=Inches(fit_h))

    add_wrapped_paragraph(s, 0.82, 3.28, 5.42, 1.32, [
        "Rectilinear perspective expands angular coverage by shrinking the focal length.",
        "As FOV approaches 180°, the projection becomes singular and boundary stretching grows rapidly."
    ], 12.3)
    add_wrapped_paragraph(s, 7.05, 3.28, 5.35, 1.32, [
        "Fisheye projection maps viewing angle to screen radius using a nonlinear function.",
        "This can show directions beyond 180° by compressing rays into a finite image."
    ], 12.3)
    add_formula(s, rendered["fov"], 1.18, 5.18, 3.45, 0.42, "rectilinear FOV scale")
    add_formula(s, rendered["theta"], 7.28, 5.15, 4.55, 0.42, "viewing angle to radial coordinate")
    add_text(s, 0.9, 6.34, 11.65, 0.35,
             "For Gaussian splats, this is not only a camera setting: centers, projected covariance, sorting, and culling must remain consistent.",
             13, True, TEAL, PP_ALIGN.CENTER)
    add_footer(s, 8)

    # Slide 2: Splat footprint.
    s = prs.slides.add_slide(blank)
    add_title(s, "Nonlinear Projection Must Transform the Splat Footprint",
              "A Gaussian splat has spatial extent, so changing the projection also changes its screen-space ellipse.")
    rect_label(s, 0.88, 1.75, 2.35, 0.72, "Meshes\nproject vertices", BLUE)
    rect_label(s, 3.85, 1.75, 2.35, 0.72, "Point clouds\nproject centers", RED)
    rect_label(s, 6.82, 1.75, 2.35, 0.72, "Gaussian splats\ncenter + footprint", TEAL)
    rect_label(s, 10.0, 1.55, 2.05, 1.08, "2D ellipse\nscreen-space covariance", TEAL, RGBColor(247, 252, 252), 10)
    arrow(s, 9.22, 2.12, 9.95, 2.12, TEAL)

    add_text(s, 1.05, 3.15, 11.1, 0.52,
             "Meshes and point clouds mainly require projecting vertices or centers. "
             "Gaussian splats also require projecting their 3D covariance into a 2D screen-space ellipse.",
             14, False, BLACK)
    add_text(s, 1.05, 3.78, 11.1, 0.52,
             "Under fisheye projection, transforming only the center is not enough: the footprint, culling, "
             "and sorting must remain consistent.",
             14, False, BLACK)
    add_formula(s, rendered["cov"], 1.65, 5.0, 4.8, 0.48, "perspective / general covariance projection")
    add_formula(s, rendered["fish_cov"], 6.85, 5.0, 4.65, 0.48, "for direct fisheye, J becomes J_fish")
    add_text(s, 1.4, 6.25, 10.4, 0.26,
             "W: view transform    ·    J: projection Jacobian    ·    Σ3D: 3D splat shape    ·    Σ2D: screen-space footprint",
             10, False, GREY, PP_ALIGN.CENTER)
    add_footer(s, 9)

    prs.save(OUT)
    strip_visual_effects_from_pptx(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
