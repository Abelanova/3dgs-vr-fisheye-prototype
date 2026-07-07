from pathlib import Path
from tempfile import NamedTemporaryFile
from zipfile import ZIP_DEFLATED, ZipFile

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"E:\3dgs fisheye projection")
SRC = ROOT / "Presentation" / "FOV_Fisheye_3DGS_slide11_12_architecture_formats.pptx"
OUT = ROOT / "Presentation" / "FOV_Fisheye_3DGS_slide12_two_small_tables.pptx"

BLACK = RGBColor(20, 24, 30)
GREY = RGBColor(82, 86, 94)
MUTED = RGBColor(118, 124, 132)
TEAL = RGBColor(0, 132, 143)
RULE = RGBColor(72, 78, 86)
LIGHT = RGBColor(214, 220, 226)
HEADER = RGBColor(241, 246, 247)


def clear_slide(slide):
    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        tree.remove(shape._element)


def fit_frame(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.025)
    tf.margin_right = Inches(0.025)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)


def text_box(slide, x, y, w, h, text, size=7, bold=False, color=BLACK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    fit_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def rule(slide, x, y, w, color=RULE, h=0.006):
    return rect(slide, x, y, w, h, color)


def table(slide, x, y, w, title, columns, rows, widths):
    text_box(slide, x, y, w, 0.18, title, 8.0, True, TEAL)
    text_box(slide, x, y + 0.26, w, 0.14, f"Table. {title.lower()}.", 5.0, False, GREY, italic=True)

    y0 = y + 0.48
    header_h = 0.36
    row_h = 0.47
    rect(slide, x, y0, w, header_h, HEADER)
    rule(slide, x, y0, w, RULE)
    cx = x
    for label, cw in zip(columns, widths):
        text_box(slide, cx + 0.03, y0 + 0.11, cw - 0.06, 0.13, label, 5.7, True)
        cx += cw
    rule(slide, x, y0 + header_h, w, RULE)

    for ri, row in enumerate(rows):
        ry = y0 + header_h + ri * row_h
        cx = x
        for ci, (value, cw) in enumerate(zip(row, widths)):
            text_box(slide, cx + 0.03, ry + 0.10, cw - 0.06, row_h - 0.12, value, 5.65, ci == 0)
            cx += cw
        if ri < len(rows) - 1:
            rule(slide, x, ry + row_h, w, LIGHT, 0.004)
    rule(slide, x, y0 + header_h + len(rows) * row_h, w, RULE)


def footer(slide):
    text_box(slide, 0.54, 5.31, 4.35, 0.14, "3DGS FOV/Fisheye Rendering for VR Scene Exploration", 4.8, False, MUTED)
    text_box(slide, 9.34, 5.31, 0.2, 0.14, "12", 5.0, False, MUTED, PP_ALIGN.RIGHT)


def strip_effects(path):
    tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)
    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(path)


def main():
    prs = Presentation(SRC)
    slide = prs.slides[11]
    clear_slide(slide)

    text_box(slide, 0.54, 0.38, 8.9, 0.36, "Gaussian Splat Formats and Tool Context", 18.5, True)
    text_box(
        slide,
        0.56,
        0.78,
        8.65,
        0.25,
        "We normalize several splat representations, then compare our Unity/VR prototype against existing 3DGS toolchains.",
        8.0,
        False,
        GREY,
    )
    rect(slide, 0.54, 1.10, 0.58, 0.03, TEAL)

    table(
        slide,
        0.68,
        1.42,
        4.10,
        "Format support",
        ["Format", "Where used", "Handling in our prototype"],
        [
            ["PLY", "Original 3DGS / common exports", "Importer + large-scene preview"],
            ["SPZ", "Compact capture-oriented splats", "UnityGaussianSplatting path"],
            ["SOG", "PlayCanvas web ecosystem", "SOG importer + color/memory fixes"],
            ["LOD", "Runtime preview representation", "Reduced budgets for VR interaction"],
        ],
        [0.58, 1.58, 1.94],
    )

    table(
        slide,
        5.20,
        1.42,
        4.12,
        "Related toolchains",
        ["Tool", "Strength", "How we use it"],
        [
            ["UnityGaussianSplatting", "Unity renderer + assets", "Base renderer"],
            ["PlayCanvas", "FOV/fisheye splat reference", "Method reference"],
            ["Nerfstudio", "Training/export ecosystem", "Pipeline context"],
        ],
        [1.28, 1.50, 1.34],
    )

    text_box(
        slide,
        0.86,
        4.78,
        8.0,
        0.20,
        "Takeaway: format support and tool comparison define the system boundary before the rendering methods are introduced.",
        6.7,
        True,
        TEAL,
        PP_ALIGN.CENTER,
    )
    footer(slide)

    prs.save(OUT)
    strip_effects(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
