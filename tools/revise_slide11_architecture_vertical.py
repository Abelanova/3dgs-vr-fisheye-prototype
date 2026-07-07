from pathlib import Path
from tempfile import NamedTemporaryFile
from zipfile import ZIP_DEFLATED, ZipFile

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"E:\3dgs fisheye projection")
SRC = ROOT / "Presentation" / "FOV_Fisheye_3DGS_clean_LaTeX_formulas_aspect.pptx"
OUT = ROOT / "Presentation" / "FOV_Fisheye_3DGS_slide11_architecture_vertical.pptx"

BLACK = RGBColor(20, 24, 30)
GREY = RGBColor(82, 86, 94)
MUTED = RGBColor(118, 124, 132)
LINE = RGBColor(166, 174, 184)
PALE = RGBColor(243, 246, 249)
MID = RGBColor(226, 231, 237)
TEAL = RGBColor(0, 132, 143)
WHITE = RGBColor(255, 255, 255)


def clear_slide(slide):
    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        tree.remove(shape._element)


def fit_frame(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)


def text_box(slide, x, y, w, h, text, size=9, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    fit_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def box(slide, x, y, w, h, text, size=7.2, bold=False, fill=WHITE, line=LINE, color=BLACK):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(0.8)
    tf = sh.text_frame
    tf.clear()
    fit_frame(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return sh


def arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = GREY
    conn.line.width = Pt(0.9)
    conn.line.end_arrowhead = True
    return conn


def footer(slide):
    text_box(slide, 0.54, 5.31, 4.35, 0.14, "3DGS FOV/Fisheye Rendering for VR Scene Exploration", 4.8, False, MUTED)
    text_box(slide, 9.34, 5.31, 0.2, 0.14, "11", 5.0, False, MUTED, PP_ALIGN.RIGHT)


def strip_effects(path):
    tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)
    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(path)


def main():
    prs = Presentation(SRC)
    slide = prs.slides[10]
    clear_slide(slide)

    text_box(slide, 0.54, 0.38, 8.9, 0.36, "System Architecture", 19, True)
    text_box(
        slide,
        0.56,
        0.78,
        8.65,
        0.34,
        "A Unity/VR prototype for evaluating nonlinear view transformations of 3D Gaussian splats through a shared asset pipeline, two rendering backends, and interactive preview interfaces.",
        8.3,
        False,
        GREY,
    )
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.54), Inches(1.12), Inches(0.58), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    # Reference-inspired vertical architecture block.
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.92), Inches(1.34), Inches(3.85), Inches(3.72))
    panel.fill.solid()
    panel.fill.fore_color.rgb = PALE
    panel.line.color.rgb = MID
    panel.line.width = Pt(0.7)

    box(slide, 1.55, 1.18, 2.58, 0.28, "3DGS scene sources", 7.0, True, fill=MID, line=MID)
    box(slide, 1.35, 1.62, 2.98, 0.42, "Asset ingestion\nPLY / SPZ / SOG import + decoding", 6.8, True)
    arrow(slide, 2.84, 2.04, 2.84, 2.22)
    box(slide, 1.35, 2.24, 2.98, 0.38, "Unity scene representation\nGaussianSplatAsset", 6.8, True)
    arrow(slide, 2.84, 2.62, 2.84, 2.80)
    box(slide, 1.35, 2.82, 2.98, 0.34, "Gaussian splat renderer core", 6.8, True)
    arrow(slide, 2.84, 3.16, 2.84, 3.34)
    box(slide, 1.18, 3.36, 1.42, 0.52, "cubemap\nbackend", 6.6)
    box(slide, 3.02, 3.36, 1.42, 0.52, "direct\nbackend", 6.6)
    arrow(slide, 2.84, 3.88, 2.84, 4.08)
    box(slide, 1.35, 4.10, 2.98, 0.42, "Interaction layer\nFOV/lens controls + desktop/VR preview", 6.6, True)
    arrow(slide, 2.84, 4.52, 2.84, 4.70)
    box(slide, 1.65, 4.72, 2.38, 0.28, "method comparison + demo output", 6.5, True, fill=MID, line=MID)

    # Side notes: function, not chronology.
    text_box(slide, 5.28, 1.48, 3.9, 0.22, "Functional decomposition", 10.2, True, TEAL)
    text_box(
        slide,
        5.30,
        1.90,
        3.8,
        0.55,
        "Asset pipeline: normalizes PLY, SPZ, and SOG scenes into Unity-side Gaussian splat assets.",
        8.2,
        False,
        BLACK,
    )
    text_box(
        slide,
        5.30,
        2.62,
        3.8,
        0.72,
        "Rendering backends: compare image-space cubemap fisheye against direct covariance-aware fisheye rendering.",
        8.2,
        False,
        BLACK,
    )
    text_box(
        slide,
        5.30,
        3.52,
        3.8,
        0.62,
        "Interactive preview: exposes the same scene and controls in desktop and VR for qualitative evaluation.",
        8.2,
        False,
        BLACK,
    )
    text_box(
        slide,
        5.30,
        4.52,
        3.8,
        0.32,
        "This structure lets us isolate failures in projection, footprint transformation, sorting, and culling.",
        7.6,
        True,
        TEAL,
    )

    footer(slide)
    prs.save(OUT)
    strip_effects(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
