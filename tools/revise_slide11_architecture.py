from pathlib import Path
from tempfile import NamedTemporaryFile
from zipfile import ZIP_DEFLATED, ZipFile

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"E:\3dgs fisheye projection")
SRC = ROOT / "Presentation" / "FOV_Fisheye_3DGS_clean_LaTeX_formulas_aspect.pptx"
OUT = ROOT / "Presentation" / "FOV_Fisheye_3DGS_slide11_architecture.pptx"

BLACK = RGBColor(20, 24, 30)
GREY = RGBColor(82, 86, 94)
MUTED = RGBColor(118, 124, 132)
LINE = RGBColor(178, 186, 196)
PALE = RGBColor(245, 247, 249)
MID = RGBColor(232, 236, 241)
TEAL = RGBColor(0, 132, 143)
WHITE = RGBColor(255, 255, 255)


def clear_slide(slide):
    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        tree.remove(shape._element)


def fit_frame(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)


def text_box(slide, x, y, w, h, text, size=9, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    fit_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def box(slide, x, y, w, h, text, size=7.5, bold=False, fill=WHITE, line=LINE, color=BLACK):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(0.8)
    tf = sh.text_frame
    tf.clear()
    fit_frame(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return sh


def panel(slide, x, y, w, h, title):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = PALE
    sh.line.color.rgb = MID
    sh.line.width = Pt(0.7)
    text_box(slide, x + 0.08, y + 0.06, w - 0.16, 0.16, title.upper(), 5.8, True, TEAL, PP_ALIGN.CENTER)
    return sh


def arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = GREY
    conn.line.width = Pt(0.9)
    conn.line.end_arrowhead = True
    return conn


def footer(slide):
    text_box(slide, 0.54, 5.31, 4.35, 0.14, "3DGS FOV/Fisheye Rendering for VR Scene Exploration", 4.8, False, MUTED)
    text_box(slide, 9.34, 5.31, 0.2, 0.14, "11", 5.0, False, MUTED, PP_ALIGN.RIGHT)


def strip_effects(path):
    tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)
    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(path)


def main():
    prs = Presentation(SRC)
    s = prs.slides[10]
    clear_slide(s)

    text_box(s, 0.54, 0.38, 8.9, 0.36, "System Architecture", 19, True)
    text_box(
        s,
        0.56,
        0.78,
        8.55,
        0.34,
        "A Unity/VR prototype that separates scene ingestion, nonlinear rendering backends, and interactive evaluation for FOV and fisheye 3D Gaussian splat exploration.",
        8.4,
        False,
        GREY,
    )
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.54), Inches(1.12), Inches(0.58), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    # Main architecture diagram: three functional layers.
    x, y, w = 0.82, 1.42, 8.35
    panel(s, x, y, w, 0.92, "Asset pipeline")
    box(s, 1.08, 1.78, 1.25, 0.32, "PLY / SPZ / SOG", 7.0, True)
    arrow(s, 2.43, 1.94, 2.78, 1.94)
    box(s, 2.92, 1.70, 1.55, 0.48, "import +\ndecode", 7.0, True)
    arrow(s, 4.58, 1.94, 4.92, 1.94)
    box(s, 5.06, 1.70, 1.72, 0.48, "Gaussian\nSplatAsset", 7.0, True)
    arrow(s, 6.90, 1.94, 7.24, 1.94)
    box(s, 7.38, 1.70, 1.25, 0.48, "Unity\nscene data", 7.0, True)

    y2 = 2.52
    panel(s, x, y2, w, 1.25, "Rendering layer")
    box(s, 1.18, 2.88, 1.72, 0.46, "3DGS renderer\nbase pass", 7.0, True)
    arrow(s, 3.02, 3.11, 3.34, 3.11)
    box(s, 3.48, 2.75, 2.0, 0.72, "cubemap backend\n6/12 perspective captures\n+ fisheye composition", 6.2)
    box(s, 5.85, 2.75, 2.25, 0.72, "direct backend\nnonlinear projection\n+ covariance Jacobian", 6.2)
    arrow(s, 5.48, 3.11, 5.78, 3.11)
    text_box(s, 8.28, 2.83, 0.55, 0.52, "sort\ncull\nblend", 6.2, True, GREY, PP_ALIGN.CENTER)

    y3 = 3.98
    panel(s, x, y3, w, 0.86, "Interaction and evaluation")
    box(s, 1.25, 4.34, 1.85, 0.34, "runtime controls\nFOV / lens / mode", 6.5)
    arrow(s, 3.22, 4.51, 3.62, 4.51)
    box(s, 3.76, 4.34, 1.5, 0.34, "desktop\npreview", 6.8, True)
    box(s, 5.58, 4.34, 1.5, 0.34, "VR / XR\npreview", 6.8, True)
    arrow(s, 7.20, 4.51, 7.58, 4.51)
    box(s, 7.72, 4.34, 1.1, 0.34, "method\ncomparison", 6.4, True)

    text_box(
        s,
        0.84,
        4.98,
        8.25,
        0.2,
        "Design focus: the same scene and controls can be evaluated across cubemap and direct fisheye rendering paths.",
        7.1,
        False,
        TEAL,
        PP_ALIGN.CENTER,
    )
    footer(s)

    prs.save(OUT)
    strip_effects(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
