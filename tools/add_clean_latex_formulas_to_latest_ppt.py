from pathlib import Path
from zipfile import ZipFile, ZIP_DEFLATED
from tempfile import NamedTemporaryFile

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import rcParams

rcParams["mathtext.fontset"] = "stix"
rcParams["font.family"] = "STIXGeneral"


SRC = Path(r"C:\Users\nova\Downloads\FOV and Fisheye Distortion Rendering of 3D Gaussian Splats for VR Scene Exploration (2).pptx")
OUT = Path(r"E:\3dgs fisheye projection\Presentation\FOV_Fisheye_3DGS_clean_LaTeX_formulas.pptx")
FORMULA_DIR = OUT.parent / "formula_renders_clean"

BLACK = RGBColor(18, 18, 18)
GREY = RGBColor(85, 85, 85)
TEAL = RGBColor(0, 135, 145)
BLUE = RGBColor(43, 93, 173)
PURPLE = RGBColor(113, 82, 165)


def render_formula(expr: str, out: Path, fontsize=25, color="#202020"):
    out.parent.mkdir(parents=True, exist_ok=True)

    # First pass: measure text extent.
    fig = plt.figure(figsize=(0.01, 0.01), dpi=320)
    text = fig.text(0, 0, expr, fontsize=fontsize, color=color)
    fig.canvas.draw()
    bbox = text.get_window_extent()
    w = max(bbox.width / fig.dpi + 0.18, 0.35)
    h = max(bbox.height / fig.dpi + 0.16, 0.25)
    plt.close(fig)

    # Second pass: render centered on transparent background.
    fig = plt.figure(figsize=(w, h), dpi=320)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")
    ax.text(0.5, 0.5, expr, ha="center", va="center", fontsize=fontsize, color=color)
    fig.savefig(out, transparent=True, bbox_inches="tight", pad_inches=0.02)
    plt.close(fig)


def add_caption(slide, x, y, w, text, color=GREY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.18))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(7.5)
    r.font.color.rgb = color
    return box


def fit_size(image_path, max_w, max_h=None):
    with Image.open(image_path) as im:
        px_w, px_h = im.size
    aspect = px_h / px_w
    if max_h is None:
        return max_w, max_w * aspect
    fit_w = max_w
    fit_h = fit_w * aspect
    if fit_h > max_h:
        fit_h = max_h
        fit_w = fit_h / aspect
    return fit_w, fit_h


def add_formula(slide, key, image_path, x, y, w, h=None, caption=None, caption_color=GREY):
    fit_w, fit_h = fit_size(image_path, w, h)
    x_offset = (w - fit_w) / 2 if w > fit_w else 0
    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(x + x_offset),
        Inches(y),
        width=Inches(fit_w),
        height=Inches(fit_h),
    )
    try:
        pic.shadow.inherit = False
    except Exception:
        pass
    if caption:
        add_caption(slide, x, y + fit_h + 0.02, w, caption, caption_color)
    return pic


def strip_visual_effects_from_pptx(path: Path):
    effect_tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)
    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in effect_tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(path)


FORMULAS = {
    # Slide 8: FOV vs fisheye
    "fov": (r"$f_y=\frac{1}{\tan(\mathrm{FOV}_y/2)}$", 27),
    "theta": (r"$\theta=\operatorname{atan2}\!\left(\sqrt{x^2+y^2},-z\right),\quad r=g(\theta)$", 24),
    "fisheye": (r"$g(\theta)=k\,\tan\!\left(\frac{\theta}{k}\right)$", 27),

    # Slide 9/10: splat projection theory
    "gaussian": (r"$G(x)=\exp\!\left[-\frac{1}{2}(x-\mu)^{T}\Sigma^{-1}(x-\mu)\right]$", 24),
    "covariance": (r"$\Sigma_{2D}=J\,W\,\Sigma_{3D}\,W^{T}J^{T}$", 27),

    # Implementation story
    "center_only": (r"$\mu\mapsto\pi_{\mathrm{fish}}(\mu)\quad\mathrm{but}\quad \Sigma_{2D}\ \mathrm{unchanged}$", 25),
    "point_debug": (r"$\pi_{\mathrm{fish}}(\mu_i)\ \mathrm{correct}\;\;\not\Rightarrow\;\;\mathrm{splat\ footprint\ correct}$", 24),

    # Cubemap method
    "cubemap_dir": (r"$d(u,v)=R_{\mathrm{cam}}\;\pi_{\mathrm{fish}}^{-1}(u,v)$", 25),
    "cubemap_sample": (r"$I(u,v)=C_{\mathrm{face}}\!\left(d(u,v)\right)$", 26),
    "vr_faces": (r"$N_{\mathrm{captures}}=6\times N_{\mathrm{eyes}}\quad(=12\ \mathrm{in\ stereo\ VR})$", 25),

    # Direct method
    "direct_chain": (r"$\mu_{\mathrm{view}}\rightarrow\theta\rightarrow\mu_{\mathrm{clip}}^{\mathrm{fish}}$", 26),
    "fish_cov": (r"$\Sigma_{2D}^{\mathrm{fish}}=J_{\mathrm{fish}}\,W\,\Sigma_{3D}\,W^{T}J_{\mathrm{fish}}^{T}$", 26),
    "sorting": (r"$\mathrm{Perspective:}\ z_{\mathrm{view}}\qquad\mathrm{Fisheye:}\ \|p_{\mathrm{view}}\|$", 25),
}


def main():
    if not SRC.exists():
        raise FileNotFoundError(SRC)
    OUT.parent.mkdir(exist_ok=True)
    FORMULA_DIR.mkdir(exist_ok=True)

    rendered = {}
    for key, (expr, size) in FORMULAS.items():
        p = FORMULA_DIR / f"{key}.png"
        render_formula(expr, p, fontsize=size)
        rendered[key] = p

    prs = Presentation(SRC)
    slides = prs.slides

    # Page 3 intentionally left unchanged.

    # Slide 8: FOV vs Fisheye
    if len(slides) >= 8:
        add_formula(slides[7], "fov", rendered["fov"], 0.95, 5.40, 3.1, 0.42,
                    "Rectilinear FOV scale")
        add_formula(slides[7], "theta", rendered["theta"], 4.35, 5.38, 4.3, 0.42,
                    "Viewing angle to radial screen coordinate")
        add_formula(slides[7], "fisheye", rendered["fisheye"], 9.05, 5.40, 2.85, 0.42,
                    "Generalized fisheye")

    # Slide 9: Why nonlinear projection is hard
    if len(slides) >= 9:
        add_formula(slides[8], "covariance", rendered["covariance"], 3.05, 5.17, 6.85, 0.46,
                    "Projected splat footprint depends on the projection Jacobian")

    # Slide 10: splats background / different from points
    if len(slides) >= 10:
        add_formula(slides[9], "gaussian", rendered["gaussian"], 1.25, 5.18, 5.7, 0.48,
                    "A splat has anisotropic shape, not only a center")
        add_formula(slides[9], "covariance", rendered["covariance"], 7.15, 5.18, 4.55, 0.48,
                    "Its projected ellipse changes with the camera model")

    # Slide 16: Attempt 1
    if len(slides) >= 16:
        add_formula(slides[15], "center_only", rendered["center_only"], 2.0, 5.35, 8.95, 0.46,
                    "First direct attempt: transform centers, but not the full footprint", TEAL)

    # Slide 18: point debugging
    if len(slides) >= 18:
        add_formula(slides[17], "point_debug", rendered["point_debug"], 2.1, 5.48, 8.85, 0.42,
                    "Debug view verifies centers; remaining problem is the splat ellipse", TEAL)

    # Slide 20: cubemap method overview
    if len(slides) >= 20:
        add_formula(slides[19], "cubemap_dir", rendered["cubemap_dir"], 7.85, 3.85, 3.85, 0.42,
                    "Screen pixel → fisheye ray direction")
        add_formula(slides[19], "cubemap_sample", rendered["cubemap_sample"], 7.95, 4.70, 3.65, 0.42,
                    "Sample the corresponding cubemap face")

    # Slide 22: VR cubemap
    if len(slides) >= 22:
        add_formula(slides[21], "vr_faces", rendered["vr_faces"], 3.65, 4.38, 5.95, 0.44,
                    "Stereo VR doubles cubemap capture work")

    # Slide 24: direct method
    if len(slides) >= 24:
        add_formula(slides[23], "direct_chain", rendered["direct_chain"], 1.15, 5.02, 4.25, 0.42,
                    "Center projection")
        add_formula(slides[23], "fish_cov", rendered["fish_cov"], 5.75, 5.02, 5.6, 0.42,
                    "Covariance-aware footprint")

    # Slide 25: FOV/Fisheye separate controls
    if len(slides) >= 25:
        add_formula(slides[24], "fov", rendered["fov"], 1.0, 5.90, 3.2, 0.40,
                    "FOV controls perspective scale")
        add_formula(slides[24], "fisheye", rendered["fisheye"], 7.55, 5.90, 3.05, 0.40,
                    "Fisheye controls nonlinear radial mapping")

    # Slide 26: covariance transform
    if len(slides) >= 26:
        add_formula(slides[25], "fish_cov", rendered["fish_cov"], 2.05, 5.05, 8.95, 0.48,
                    "Final direct method transforms the projected covariance with the fisheye Jacobian")

    # Slide 27: sorting
    if len(slides) >= 27:
        add_formula(slides[26], "sorting", rendered["sorting"], 6.95, 5.00, 4.9, 0.46,
                    "Wide/fisheye views benefit from radial distance sorting")

    prs.save(OUT)
    strip_visual_effects_from_pptx(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
