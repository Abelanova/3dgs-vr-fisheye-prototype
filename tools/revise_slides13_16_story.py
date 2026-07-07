from io import BytesIO
from pathlib import Path
from tempfile import NamedTemporaryFile
from zipfile import ZIP_DEFLATED, ZipFile

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"E:\3dgs fisheye projection")
SRC = ROOT / "Presentation" / "FOV_Fisheye_3DGS_slide12_academic_tables.pptx"
OUT = ROOT / "Presentation" / "FOV_Fisheye_3DGS_slides13_16_story.pptx"

ARTIFACT_IMG = Path(r"C:\Users\nova\AppData\Local\Temp\codex-clipboard-3e0de5ba-5129-4d68-9d99-b3a0dbcc8fb0.png")
POINT_DEBUG_IMG = Path(r"C:\Users\nova\AppData\Local\Temp\codex-clipboard-5ca3395c-9db6-4f46-befa-f0ceb8ba4bac.png")

BLACK = RGBColor(20, 24, 30)
GREY = RGBColor(82, 86, 94)
MUTED = RGBColor(118, 124, 132)
TEAL = RGBColor(0, 132, 143)
RED = RGBColor(170, 55, 55)
RULE = RGBColor(74, 82, 92)
LIGHT = RGBColor(222, 226, 231)
PALE = RGBColor(246, 248, 250)
WHITE = RGBColor(255, 255, 255)


def clear_slide(slide):
    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        tree.remove(shape._element)


def fit_frame(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.035)
    tf.margin_right = Inches(0.035)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)


def text_box(slide, x, y, w, h, text, size=8, bold=False, color=BLACK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    fit_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def bullet_box(slide, x, y, w, h, items, size=7.2, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    fit_frame(tf)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = item
        r.font.name = "Aptos"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def title(slide, heading, subtitle, n):
    text_box(slide, 0.54, 0.38, 8.9, 0.36, heading, 18.5, True)
    text_box(slide, 0.56, 0.78, 8.65, 0.26, subtitle, 8.0, False, GREY)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.54), Inches(1.10), Inches(0.58), Inches(0.03))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    footer(slide, n)


def footer(slide, n):
    text_box(slide, 0.54, 5.31, 4.35, 0.14, "3DGS FOV/Fisheye Rendering for VR Scene Exploration", 4.8, False, MUTED)
    text_box(slide, 9.34, 5.31, 0.2, 0.14, str(n), 5.0, False, MUTED, PP_ALIGN.RIGHT)


def rect(slide, x, y, w, h, fill=WHITE, line=LIGHT):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(0.8)
    return sh


def plain_rule(slide, x, y, w, color=RULE, h=0.006):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def label(slide, x, y, w, h, text, size=7.2, color=BLACK):
    sh = rect(slide, x, y, w, h, WHITE, LIGHT)
    tf = sh.text_frame
    tf.clear()
    fit_frame(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = color
    return sh


def arrow(slide, x1, y1, x2, y2, color=GREY):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(0.9)
    conn.line.end_arrowhead = True
    return conn


def fit_picture(slide, path_or_blob, x, y, max_w, max_h):
    pic = slide.shapes.add_picture(path_or_blob, Inches(x), Inches(y))
    aspect = pic.width / pic.height if pic.height else 1
    w = max_w
    h = w / aspect
    if h > max_h:
        h = max_h
        w = h * aspect
    pic.width = Inches(w)
    pic.height = Inches(h)
    return pic


def image_blobs(slide):
    blobs = []
    for shape in slide.shapes:
        if hasattr(shape, "image"):
            blobs.append(shape.image.blob)
    return blobs


def strip_effects(path):
    tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)
    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(path)


def slide13(slide):
    clear_slide(slide)
    title(
        slide,
        "Attempt 1: Direct Fisheye Center Projection",
        "The first desktop prototype applied a fisheye mapping to splat centers, but left the splat footprint model mostly unchanged.",
        13,
    )
    fit_picture(slide, str(ARTIFACT_IMG), 0.70, 1.42, 5.15, 2.38)
    text_box(slide, 0.70, 3.88, 5.15, 0.18, "Observed failure: stretched or surrounding splats near wide-FOV image boundaries.", 6.4, False, GREY, PP_ALIGN.CENTER)

    text_box(slide, 6.20, 1.48, 2.85, 0.18, "What worked", 8.0, True, TEAL)
    bullet_box(slide, 6.20, 1.78, 2.90, 0.72, [
        "Desktop preview scene and runtime controls.",
        "Wider angular coverage became visible.",
        "Center locations were plausibly mapped."
    ], 6.9)
    text_box(slide, 6.20, 2.72, 2.85, 0.18, "What failed", 8.0, True, RED)
    bullet_box(slide, 6.20, 3.02, 2.90, 0.72, [
        "Large boundary artifacts remained.",
        "Footprints were not transformed consistently.",
        "Sorting/culling assumptions still matched a narrow perspective view."
    ], 6.9)
    text_box(slide, 0.82, 4.72, 8.1, 0.22, "Takeaway: changing the center projection alone is insufficient for Gaussian splats.", 7.2, True, TEAL, PP_ALIGN.CENTER)


def slide14(slide):
    clear_slide(slide)
    title(
        slide,
        "Failure Analysis: A Splat Is Not Just a Point",
        "A projected Gaussian has both a center and a screen-space ellipse; nonlinear projection must keep both consistent.",
        14,
    )
    label(slide, 0.95, 1.55, 1.40, 0.46, "center\nprojection", 7.0)
    arrow(slide, 2.48, 1.78, 3.00, 1.78)
    label(slide, 3.15, 1.55, 1.65, 0.46, "footprint\nprojection", 7.0)
    arrow(slide, 4.94, 1.78, 5.46, 1.78)
    label(slide, 5.60, 1.55, 1.40, 0.46, "sorting /\nculling", 7.0)
    arrow(slide, 7.14, 1.78, 7.66, 1.78)
    label(slide, 7.80, 1.55, 1.10, 0.46, "blend", 7.0)

    text_box(slide, 1.10, 2.44, 3.55, 0.20, "Perspective splat rendering", 8.0, True, BLACK)
    bullet_box(slide, 1.10, 2.78, 3.55, 0.82, [
        "The renderer projects a 3D covariance into a 2D ellipse.",
        "The usual footprint assumes a local perspective projection.",
        "Depth ordering is mainly forward-facing."
    ], 6.9)
    text_box(slide, 5.25, 2.44, 3.55, 0.20, "Why fisheye breaks this assumption", 8.0, True, BLACK)
    bullet_box(slide, 5.25, 2.78, 3.55, 0.82, [
        "The mapping is nonlinear over the image.",
        "Side/near-surround splats can enter the view.",
        "A very large footprint can dominate the frame."
    ], 6.9)

    text_box(slide, 2.05, 4.08, 5.90, 0.23, "Σ₂D = J W Σ₃D Wᵀ Jᵀ", 14.0, False, BLACK, PP_ALIGN.CENTER)
    text_box(slide, 2.05, 4.43, 5.90, 0.18, "The projection Jacobian J changes when the view model changes.", 6.4, False, GREY, PP_ALIGN.CENTER)
    text_box(slide, 0.92, 4.84, 8.05, 0.18, "Diagnosis: the artifact is a footprint/covariance problem, not only a camera-FOV problem.", 7.0, True, TEAL, PP_ALIGN.CENTER)


def slide15(slide):
    clear_slide(slide)
    title(
        slide,
        "Diagnostic View: Point Rendering",
        "We used the renderer's debug point mode to isolate center projection from splat footprint rendering.",
        15,
    )
    fit_picture(slide, str(POINT_DEBUG_IMG), 0.70, 1.36, 5.20, 3.18)
    text_box(slide, 0.70, 4.60, 5.20, 0.16, "Point/debug view: only centers are drawn, so the ellipse footprint is bypassed.", 6.2, False, GREY, PP_ALIGN.CENTER)
    text_box(slide, 6.20, 1.42, 2.85, 0.18, "What this verified", 8.0, True, TEAL)
    bullet_box(slide, 6.20, 1.74, 2.90, 0.92, [
        "Fisheye-mapped centers remain spatially coherent.",
        "The surrounding-splat artifact does not appear in point mode.",
        "The failure therefore comes after center projection."
    ], 6.9)
    text_box(slide, 6.20, 2.96, 2.85, 0.18, "Next implication", 8.0, True, BLACK)
    bullet_box(slide, 6.20, 3.27, 2.90, 0.74, [
        "Either post-distort a robust perspective rendering,",
        "or modify the splat renderer to transform footprints directly."
    ], 6.9)
    text_box(slide, 0.90, 4.94, 8.05, 0.18, "Takeaway: point rendering ruled out center placement as the main source of the artifact.", 7.0, True, TEAL, PP_ALIGN.CENTER)


def slide16(slide, preserved_blobs):
    clear_slide(slide)
    title(
        slide,
        "Attempt 2: Cubemap Fisheye Baseline",
        "The second approach renders ordinary perspective cubemap faces and applies the fisheye mapping only in a composite pass.",
        16,
    )
    for i, name in enumerate(["+X", "-X", "+Y", "-Y", "+Z", "-Z"]):
        label(slide, 0.78 + i * 0.74, 1.46, 0.48, 0.32, name, 6.4)
    text_box(slide, 0.78, 1.93, 4.18, 0.16, "six 90° perspective captures", 6.2, False, GREY, PP_ALIGN.CENTER)
    arrow(slide, 5.10, 1.62, 5.68, 1.62)
    label(slide, 5.84, 1.36, 1.45, 0.54, "cubemap\ntextures", 6.6)
    arrow(slide, 7.42, 1.62, 7.96, 1.62)
    label(slide, 8.10, 1.36, 1.20, 0.54, "fisheye\ncomposite", 6.6)

    text_box(slide, 0.85, 2.44, 3.75, 0.18, "Composite mapping", 8.0, True, BLACK)
    bullet_box(slide, 0.85, 2.74, 3.75, 0.88, [
        "For each output pixel, compute a fisheye ray direction.",
        "Sample the matching cubemap face.",
        "The splat renderer itself remains perspective per face."
    ], 6.9)
    text_box(slide, 5.22, 2.44, 3.75, 0.18, "Why this helped", 8.0, True, BLACK)
    bullet_box(slide, 5.22, 2.74, 3.75, 0.88, [
        "Provides a robust visual baseline.",
        "Avoids direct nonlinear footprint projection at first.",
        "Useful for side-by-side comparison with direct rendering."
    ], 6.9)

    text_box(slide, 1.30, 4.02, 7.40, 0.20, "screen pixel → fisheye ray direction → cubemap face sample", 8.0, False, BLACK, PP_ALIGN.CENTER)
    plain_rule(slide, 1.60, 4.34, 6.80, LIGHT, 0.004)
    text_box(slide, 1.10, 4.64, 7.90, 0.18, "Trade-off: conceptually stable, but expensive for VR because stereo requires many perspective captures.", 7.0, True, TEAL, PP_ALIGN.CENTER)

    # Preserve the existing source/diagram images as a compact reference strip.
    x = 7.25
    for idx, blob in enumerate(preserved_blobs[:2]):
        fit_picture(slide, BytesIO(blob), x + idx * 0.82, 4.05, 0.72, 0.42)
    if preserved_blobs:
        text_box(slide, 7.18, 4.50, 1.85, 0.14, "existing formula/source assets kept", 4.7, False, MUTED, PP_ALIGN.CENTER)


def main():
    prs = Presentation(SRC)
    preserved = image_blobs(prs.slides[15])
    slide13(prs.slides[12])
    slide14(prs.slides[13])
    slide15(prs.slides[14])
    slide16(prs.slides[15], preserved)
    prs.save(OUT)
    strip_effects(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
