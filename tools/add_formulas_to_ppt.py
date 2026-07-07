from pathlib import Path
from zipfile import ZipFile, ZIP_DEFLATED
from tempfile import NamedTemporaryFile
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


SRC = Path(r"C:\Users\nova\Downloads\FOV and Fisheye Distortion Rendering of 3D Gaussian Splats for VR Scene Exploration (1).pptx")
OUT = Path(r"E:\3dgs fisheye projection\Presentation\FOV_Fisheye_3DGS_with_formulas.pptx")

BLACK = RGBColor(20, 20, 20)
GREY = RGBColor(85, 85, 85)
TEAL = RGBColor(0, 135, 145)
BLUE = RGBColor(43, 93, 173)
PURPLE = RGBColor(113, 82, 165)
LIGHT_TEAL = RGBColor(237, 250, 251)
LIGHT_BLUE = RGBColor(238, 244, 252)
LIGHT_PURPLE = RGBColor(246, 242, 251)
LIGHT_GREY = RGBColor(248, 249, 250)
MID = RGBColor(214, 222, 230)


def add_formula_box(slide, x, y, w, h, title, formula_lines, note=None, accent=TEAL, fill=LIGHT_TEAL):
    """Add a flat academic formula box. Formula lines stay editable."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = accent
    box.line.width = Pt(1.1)
    try:
        box.shadow.inherit = False
    except Exception:
        pass

    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos"
    r.font.bold = True
    r.font.size = Pt(10)
    r.font.color.rgb = accent

    for line in formula_lines:
        p = tf.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(0)
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.name = "Cambria Math"
        r.font.size = Pt(16)
        r.font.color.rgb = BLACK

    if note:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = note
        r.font.name = "Aptos"
        r.font.size = Pt(8.5)
        r.font.color.rgb = GREY

    return box


def add_small_note(slide, x, y, w, h, text, accent=TEAL):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(9)
    r.font.color.rgb = accent
    r.font.bold = True
    return box


def strip_visual_effects_from_pptx(path: Path):
    effect_tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)

    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in effect_tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(path)


def main():
    if not SRC.exists():
        raise FileNotFoundError(SRC)
    OUT.parent.mkdir(exist_ok=True)

    prs = Presentation(SRC)
    slides = prs.slides

    # Slide numbers are 1-based in the user's PowerPoint.
    # 3: What Is 3D Gaussian Splatting?
    if len(slides) >= 3:
        add_formula_box(
            slides[2], 7.05, 4.65, 4.95, 1.15,
            "3D Gaussian representation",
            ["G(x) = exp(-1/2 (x - μ)^T Σ^{-1} (x - μ))"],
            "μ is the splat center; Σ encodes anisotropic 3D shape.",
            TEAL, LIGHT_TEAL
        )

    # 8: FOV vs Fisheye
    if len(slides) >= 8:
        add_formula_box(
            slides[7], 0.92, 5.33, 5.35, 1.15,
            "Rectilinear FOV",
            ["f_y = 1 / tan(FOV_y / 2)"],
            "As FOV approaches 180°, perspective projection stretches strongly near the boundary.",
            BLUE, LIGHT_BLUE
        )
        add_formula_box(
            slides[7], 6.92, 5.33, 5.35, 1.15,
            "Fisheye mapping",
            ["r = g(θ),     θ = atan2(√(x²+y²), -z)"],
            "Fisheye maps viewing angle θ nonlinearly instead of using a flat image plane.",
            PURPLE, LIGHT_PURPLE
        )

    # 9: Why nonlinear projection is hard
    if len(slides) >= 9:
        add_formula_box(
            slides[8], 1.15, 5.15, 10.65, 0.95,
            "Projected covariance is part of rendering",
            ["Σ_2D = J W Σ_3D W^T J^T"],
            "W transforms the Gaussian into camera space; J is the projection Jacobian.",
            TEAL, LIGHT_TEAL
        )

    # 15: PlayCanvas reference
    if len(slides) >= 15:
        add_formula_box(
            slides[14], 7.02, 3.05, 4.88, 1.1,
            "Generalized fisheye model",
            ["g(θ) = k · tan(θ / k)"],
            "This is the model we use as a reference for direct fisheye projection.",
            PURPLE, LIGHT_PURPLE
        )

    # 17: problem at edges
    if len(slides) >= 17:
        add_formula_box(
            slides[16], 0.98, 5.33, 10.95, 0.92,
            "Why center-only projection fails",
            ["μ → π_fish(μ)  is not enough;     Σ_2D also needs J_fish"],
            "The splat center can be correct while the ellipse footprint is still wrong.",
            TEAL, LIGHT_TEAL
        )

    # 19: lesson
    if len(slides) >= 19:
        add_formula_box(
            slides[18], 1.1, 5.05, 10.7, 0.95,
            "Renderer-level requirement",
            ["center projection + covariance projection + sorting + culling"],
            "A fisheye splat renderer must keep all four parts consistent.",
            BLUE, LIGHT_BLUE
        )

    # 20: cubemap
    if len(slides) >= 20:
        add_formula_box(
            slides[19], 8.75, 4.95, 3.25, 0.95,
            "Cubemap sampling",
            ["color(u,v) = C_face(d(u,v))"],
            "Distortion happens after ordinary perspective captures.",
            BLUE, LIGHT_BLUE
        )

    # 24: direct method
    if len(slides) >= 24:
        add_formula_box(
            slides[23], 0.98, 5.05, 10.95, 0.98,
            "Direct fisheye splat projection",
            ["μ_view → θ → μ_clip,fish", "Σ_2D,fish = J_fish W Σ_3D W^T J_fish^T"],
            "The direct path changes the splat renderer instead of post-distorting a rendered image.",
            TEAL, LIGHT_TEAL
        )

    # 25: projection model
    if len(slides) >= 25:
        add_formula_box(
            slides[24], 1.0, 5.93, 5.25, 0.78,
            "FOV controls scale",
            ["f_y = 1 / tan(FOV_y / 2)"],
            None,
            BLUE, LIGHT_BLUE
        )
        add_formula_box(
            slides[24], 6.95, 5.93, 5.25, 0.78,
            "Fisheye controls radial mapping",
            ["r = k · tan(θ / k)"],
            None,
            PURPLE, LIGHT_PURPLE
        )

    # 26: covariance transform
    if len(slides) >= 26:
        add_formula_box(
            slides[25], 1.35, 5.1, 10.35, 0.95,
            "Covariance transform",
            ["Σ'_2D = J_fish W Σ_3D W^T J_fish^T"],
            "This is the main theoretical reason the final direct method differs from center-only fisheye.",
            TEAL, LIGHT_TEAL
        )

    # 27: sorting
    if len(slides) >= 27:
        add_formula_box(
            slides[26], 7.05, 5.05, 4.75, 1.0,
            "Sorting heuristic",
            ["Perspective: sort by z_view", "Fisheye: sort by ||p_view||"],
            "Radial sorting is more stable when the visible view is very wide.",
            PURPLE, LIGHT_PURPLE
        )

    # 29: comparison plan
    if len(slides) >= 29:
        add_formula_box(
            slides[28], 1.0, 5.9, 10.75, 0.72,
            "Evaluation quantities",
            ["quality artifacts, FPS/GPU time, memory, latency, VR comfort"],
            None,
            TEAL, LIGHT_TEAL
        )

    prs.save(OUT)
    strip_visual_effects_from_pptx(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
