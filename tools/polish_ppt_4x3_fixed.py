from pathlib import Path
from tempfile import NamedTemporaryFile
from zipfile import ZIP_DEFLATED, ZipFile

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"E:\3dgs fisheye projection")
SRC = ROOT / "Presentation" / "FOV_Fisheye_3DGS_clean_LaTeX_formulas_aspect.pptx"
OUT = ROOT / "Presentation" / "FOV_Fisheye_3DGS_conservative_4x3_fixed.pptx"

BLACK = RGBColor(18, 18, 18)
GREY = RGBColor(82, 82, 82)
MUTED = RGBColor(120, 120, 120)
TEAL = RGBColor(0, 135, 145)
BLUE = RGBColor(43, 93, 173)
RED = RGBColor(196, 65, 62)
PURPLE = RGBColor(113, 82, 165)
GREEN = RGBColor(64, 145, 98)
GOLD = RGBColor(186, 142, 45)
WHITE = RGBColor(255, 255, 255)


def remove_nonpictures(slide):
    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        if not hasattr(shape, "image"):
            tree.remove(shape._element)


def clear_slide(slide):
    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        tree.remove(shape._element)


def fit_text_frame(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)


def text_box(slide, x, y, w, h, text, size=10, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    fit_text_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def bullet_box(slide, x, y, w, h, items, size=9.5, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    fit_text_frame(tf)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = item
        r.font.name = "Aptos"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def title(slide, heading, subtitle=None, section=None):
    if section:
        text_box(slide, 0.55, 0.25, 1.8, 0.2, section.upper(), 7.5, True, TEAL)
    text_box(slide, 0.55, 0.52, 8.75, 0.48, heading, 20.5, True)
    if subtitle:
        text_box(slide, 0.56, 1.08, 8.55, 0.32, subtitle, 9.2, False, GREY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.46), Inches(0.62), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def footer(slide, n):
    text_box(slide, 0.55, 7.08, 4.5, 0.18, "3DGS FOV/Fisheye Rendering for VR Scene Exploration", 6.2, False, MUTED)
    text_box(slide, 9.25, 7.08, 0.25, 0.18, str(n), 6.8, False, MUTED, PP_ALIGN.RIGHT)


def source(slide, text):
    text_box(slide, 0.55, 6.78, 8.85, 0.22, text, 5.8, False, MUTED)


def label(slide, x, y, w, h, text, line=TEAL, size=8.4):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = line
    sh.line.width = Pt(1)
    tf = sh.text_frame
    tf.clear()
    fit_text_frame(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = BLACK
    return sh


def arrow(slide, x1, y1, x2, y2, color=GREY):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.1)
    conn.line.end_arrowhead = True
    return conn


def pictures(slide):
    return [s for s in slide.shapes if hasattr(s, "image")]


def set_pic(shape, x, y, w=None, h=None):
    shape.left = Inches(x)
    shape.top = Inches(y)
    if w is not None:
        shape.width = Inches(w)
    if h is not None:
        shape.height = Inches(h)


def strip_effects(path):
    tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)
    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(path)


def polish():
    prs = Presentation(SRC)

    # 1. Title
    s = prs.slides[0]
    remove_nonpictures(s)
    for pic, (x, y, w) in zip(pictures(s), [(0.8, 5.7, 0.72), (1.9, 5.7, 0.72), (3.2, 5.48, 0.82), (4.35, 5.48, 0.82), (8.15, 5.42, 0.82)]):
        set_pic(pic, x, y, w=w)
    text_box(s, 0.62, 0.72, 8.65, 1.05, "FOV and Fisheye Distortion Rendering\nof 3D Gaussian Splats for VR Scene Exploration", 23.5, True)
    text_box(s, 0.64, 2.02, 8.25, 0.28, "A Unity prototype for nonlinear view transformations of Gaussian splat scenes", 11.5, False, GREY)
    for i, (txt, col) in enumerate([("PLY / SPZ / SOG", BLUE), ("Unity + URP", TEAL), ("Desktop + VR", GREEN), ("Cubemap vs Direct", PURPLE)]):
        label(s, 0.64 + i * 2.1, 2.68, 1.58, 0.32, txt, col, 8.0)
    text_box(s, 0.64, 6.78, 5.1, 0.24, "Yingfangzhong SUN  /  Supervisor: Daniel Filonik", 9.2, False)

    # 2. Roadmap
    s = prs.slides[1]
    clear_slide(s)
    title(s, "Talk Roadmap", "From 3DGS representation to nonlinear rendering in Unity/VR.")
    rows = [
        ("1", "Background", "3DGS representation and splat footprints."),
        ("2", "Motivation", "Why FOV/fisheye helps VR scene exploration."),
        ("3", "Formats", "PLY, SPZ, and SOG in a Unity asset pipeline."),
        ("4", "Implementation", "Direct attempt, point debug, cubemap, direct covariance."),
        ("5", "Comparison", "FOV vs fisheye; cubemap vs direct; desktop vs VR."),
        ("6", "Discussion", "Current status, next steps, and references."),
    ]
    for i, (num, head, desc) in enumerate(rows):
        y = 1.75 + i * 0.7
        text_box(s, 0.75, y, 0.3, 0.24, num, 13, True, [BLUE, TEAL, PURPLE, GOLD, GREEN, RED][i], PP_ALIGN.RIGHT)
        text_box(s, 1.32, y, 1.65, 0.24, head, 11.8, True)
        text_box(s, 3.35, y + 0.02, 5.7, 0.24, desc, 9.4, False, GREY)
    footer(s, 2)

    # 3. Light text cleanup, preserving the existing diagram.
    s = prs.slides[2]
    for sh in list(s.shapes):
        if hasattr(sh, "text") and sh.text.strip() and ("Core idea" in sh.text or "3DGS is" in sh.text or "input:" in sh.text or "References:" in sh.text):
            sh.text_frame.clear()
    text_box(s, 0.62, 1.38, 2.2, 0.22, "Core idea", 11.5, True, TEAL)
    bullet_box(s, 0.62, 1.8, 2.35, 1.6, ["Input: posed images.", "Representation: anisotropic 3D Gaussians.", "Rendering: project, sort, alpha-blend."], 9.3)
    text_box(s, 0.62, 3.72, 2.35, 0.42, "A splat has position, opacity, color, and shape.", 8.8, True, TEAL, PP_ALIGN.CENTER)
    source(s, "References: Kerbl et al., 3D Gaussian Splatting, SIGGRAPH 2023; Mildenhall et al., NeRF, ECCV 2020.")

    # 4. Background image + concise theory.
    s = prs.slides[3]
    remove_nonpictures(s)
    if pictures(s):
        set_pic(pictures(s)[0], 0.7, 1.78, w=5.45)
    title(s, "From Captured Images to a Renderable Scene", "3DGS learns an explicit splat representation that can be rendered from new viewpoints.", "Background")
    bullet_box(s, 6.45, 1.9, 2.9, 1.55, ["Learns radiance from posed images.", "Renders explicit splats interactively.", "Good fit for desktop and VR inspection."], 9.2)
    text_box(s, 6.45, 4.05, 2.9, 0.72, "But the camera model also changes the splat footprint.", 11.2, True, TEAL)
    source(s, "Image source: LearnOpenCV 3D Gaussian Splatting explanation. Method reference: Kerbl et al., 2023.")
    footer(s, 4)

    # 5. VR motivation.
    s = prs.slides[4]
    remove_nonpictures(s)
    if pictures(s):
        set_pic(pictures(s)[0], 0.7, 1.78, w=5.95)
    title(s, "Why 3DGS Is Interesting for VR", "Photorealistic captured spaces can be inspected interactively, but the user still has limited instantaneous view.", "Background")
    bullet_box(s, 6.95, 1.92, 2.45, 2.35, ["Real-time novel views.", "Detail comes from captured spaces.", "Useful for immersive inspection.", "Question: how should wide FOV and fisheye work for splats?"], 9.0)
    source(s, "Image source: 3DGS web viewer screenshot from current deck. Verify original URL before final submission.")
    footer(s, 5)

    # 6. Motivation.
    s = prs.slides[5]
    remove_nonpictures(s)
    title(s, "Motivation: Bring More Scene Content Into View", "FOV and fisheye controls can reveal peripheral or hidden content during VR scene exploration.", "Motivation")
    bullet_box(s, 0.68, 1.8, 3.55, 1.95, ["VR has limited instantaneous field of view.", "Relevant content may lie outside the current view.", "Nonlinear viewing can compress more context into the display."], 10.2)
    label(s, 5.0, 1.9, 1.45, 0.55, "normal view", BLUE, 8.6)
    arrow(s, 6.58, 2.18, 7.05, 2.18, TEAL)
    label(s, 7.22, 1.9, 1.75, 0.55, "wide / fisheye view", TEAL, 8.6)
    text_box(s, 4.8, 3.28, 4.35, 0.72, "Goal: interactive nonlinear viewing transformations for Gaussian splat scenes in Unity and VR.", 11.5, True, TEAL, PP_ALIGN.CENTER)
    footer(s, 6)

    # 7. Advisor demo inspirations.
    s = prs.slides[6]
    remove_nonpictures(s)
    if len(pictures(s)) >= 2:
        set_pic(pictures(s)[0], 0.65, 1.72, w=4.1)
        set_pic(pictures(s)[1], 5.2, 1.72, w=4.1)
    title(s, "Inspiration: Flexible Spatial Transformations", "We first studied simpler demos, then asked how similar transformations behave for Gaussian splats.", "Motivation")
    text_box(s, 0.8, 4.55, 3.8, 0.24, "World bending / spatial deformation demo", 9.2, True, BLACK, PP_ALIGN.CENTER)
    text_box(s, 5.35, 4.55, 3.8, 0.24, "FOV / fisheye projection demo", 9.2, True, BLACK, PP_ALIGN.CENTER)
    text_box(s, 0.9, 5.25, 8.1, 0.55, "Challenge for 3DGS: centers and screen-space footprints must transform consistently.", 11.0, True, TEAL, PP_ALIGN.CENTER)
    source(s, "Demo sources: advisor-provided fisheye and inception/world-bending demos; screenshots from current deck.")
    footer(s, 7)

    # 10. Splat projection theory.
    s = prs.slides[9]
    remove_nonpictures(s)
    pics = pictures(s)
    if pics:
        set_pic(pics[0], 0.72, 1.75, w=3.25)
    title(s, "Splats Are Harder Than Points Under Nonlinear Projection", "A Gaussian splat has a screen-space footprint, so changing projection affects more than its center.", "Background")
    label(s, 4.35, 1.75, 1.18, 0.5, "mesh\nvertices", BLUE, 8.0)
    label(s, 5.85, 1.75, 1.18, 0.5, "points\ncenters", RED, 8.0)
    label(s, 7.35, 1.75, 1.65, 0.5, "splats\ncenter + ellipse", TEAL, 7.8)
    bullet_box(s, 4.35, 2.62, 4.7, 0.95, ["Meshes/points: project vertices or centers.", "Splats: also project 3D covariance to a 2D ellipse."], 9.5)
    if len(pics) > 1:
        set_pic(pics[1], 4.28, 4.12, w=3.0)
    if len(pics) > 2:
        set_pic(pics[2], 7.48, 4.12, w=1.85)
    text_box(s, 4.35, 5.05, 4.65, 0.45, "J is the projection Jacobian. Under fisheye, J becomes the fisheye Jacobian.", 8.4, False, GREY, PP_ALIGN.CENTER)
    source(s, "Formula: standard 3DGS covariance projection; see Kerbl et al., 2023. Visual Gaussian illustration from current deck.")
    footer(s, 10)

    # 11. Project architecture.
    s = prs.slides[10]
    clear_slide(s)
    title(s, "Project Architecture", "The prototype connects asset conversion, renderer modifications, and desktop/VR preview scenes.", "Our Project")
    for i, (txt, col) in enumerate([("PLY", BLUE), ("SPZ", BLUE), ("SOG", PURPLE)]):
        label(s, 0.65, 1.72 + i * 0.42, 0.72, 0.28, txt, col, 7.6)
    text_box(s, 0.5, 1.36, 1.05, 0.2, "input formats", 7.8, True, BLACK, PP_ALIGN.CENTER)
    arrow(s, 1.5, 2.08, 1.88, 2.08)
    label(s, 2.02, 1.72, 1.25, 0.72, "importer /\nconverter", TEAL, 8.0)
    arrow(s, 3.38, 2.08, 3.78, 2.08)
    label(s, 3.9, 1.72, 1.35, 0.72, "Gaussian\nSplatAsset", TEAL, 8.0)
    arrow(s, 5.36, 2.08, 5.76, 2.08)
    label(s, 5.88, 1.52, 1.65, 0.5, "UnityGaussianSplatting", GREEN, 7.4)
    label(s, 5.88, 2.22, 1.65, 0.5, "FOV / fisheye params", GOLD, 7.4)
    arrow(s, 7.65, 2.08, 8.02, 2.08)
    label(s, 8.12, 1.55, 1.25, 0.46, "desktop\npreview", BLUE, 7.5)
    label(s, 8.12, 2.25, 1.25, 0.46, "VR / XR\npreview", PURPLE, 7.5)
    text_box(s, 0.65, 4.0, 1.75, 0.24, "implementation paths", 10.0, True)
    label(s, 0.65, 4.45, 1.8, 0.58, "center-only direct\nfirst attempt", RED, 7.2)
    label(s, 2.75, 4.45, 1.8, 0.58, "point debug\ncenter validation", GOLD, 7.2)
    label(s, 4.85, 4.45, 1.8, 0.58, "cubemap fisheye\nbaseline", BLUE, 7.2)
    label(s, 6.95, 4.45, 2.15, 0.58, "direct covariance-aware\nfisheye", TEAL, 7.0)
    text_box(s, 0.8, 5.72, 8.25, 0.42, "The presentation follows this story: what broke, how we debugged it, and how the final methods compare.", 9.4, False, GREY, PP_ALIGN.CENTER)
    footer(s, 11)

    prs.save(OUT)
    strip_effects(OUT)
    print(OUT)


if __name__ == "__main__":
    polish()
