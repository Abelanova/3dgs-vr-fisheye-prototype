from pathlib import Path
from tempfile import NamedTemporaryFile
from zipfile import ZIP_DEFLATED, ZipFile

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"E:\3dgs fisheye projection")
SRC = ROOT / "Presentation" / "FOV_Fisheye_3DGS_slide11_architecture_vertical.pptx"
OUT = ROOT / "Presentation" / "FOV_Fisheye_3DGS_slide11_12_architecture_formats.pptx"

BLACK = RGBColor(20, 24, 30)
GREY = RGBColor(82, 86, 94)
MUTED = RGBColor(118, 124, 132)
TEAL = RGBColor(0, 132, 143)
RULE = RGBColor(70, 76, 84)
LIGHT_RULE = RGBColor(215, 220, 226)
WHITE = RGBColor(255, 255, 255)


def clear_slide(slide):
    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        tree.remove(shape._element)


def fit_frame(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)


def text_box(slide, x, y, w, h, text, size=8, bold=False, color=BLACK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    fit_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def rule(slide, x, y, w, color=RULE, width=0.9):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.006))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def footer(slide):
    text_box(slide, 0.54, 5.31, 4.35, 0.14, "3DGS FOV/Fisheye Rendering for VR Scene Exploration", 4.8, False, MUTED)
    text_box(slide, 9.34, 5.31, 0.2, 0.14, "12", 5.0, False, MUTED, PP_ALIGN.RIGHT)


def strip_effects(path):
    tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)
    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(path)


def main():
    prs = Presentation(SRC)
    slide = prs.slides[11]
    clear_slide(slide)

    text_box(slide, 0.54, 0.38, 8.9, 0.36, "Gaussian Splat Data Formats", 19, True)
    text_box(
        slide,
        0.56,
        0.78,
        8.55,
        0.30,
        "The asset pipeline normalizes several 3DGS representations before they enter the same Unity rendering and preview system.",
        8.4,
        False,
        GREY,
    )
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.54), Inches(1.12), Inches(0.58), Inches(0.03))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    text_box(slide, 0.64, 1.34, 8.55, 0.17, "Table 1. Representations handled by the prototype asset pipeline.", 6.6, False, GREY, italic=True)

    # Booktabs-style table: no vertical grid, no colored blocks.
    x0 = 0.64
    table_w = 8.72
    y0 = 1.62
    col_x = [x0, 1.62, 3.68, 6.05]
    col_w = [0.82, 1.80, 2.10, 3.18]
    headers = ["Format", "Primary role", "Encoding / structure", "Prototype handling"]
    rows = [
        (
            "PLY",
            "Canonical 3DGS interchange and common training export.",
            "Per-splat attributes: position, scale, rotation, opacity, and color/SH coefficients.",
            "Imported into Unity GaussianSplatAsset; preview budgets help with large scenes.",
        ),
        (
            "SPZ",
            "Compact binary exchange format used by mobile / capture-oriented workflows.",
            "Quantized compressed splat parameters with format-specific scale and color conventions.",
            "Supported through the UnityGaussianSplatting import path and conversion checks.",
        ),
        (
            "SOG",
            "Web-oriented container used in the PlayCanvas Gaussian splat ecosystem.",
            "ZIP package with metadata, WebP payloads, quantization, and codebook-style decoding.",
            "Integrated as a Unity importer; fixed color decoding and memory issues during loading.",
        ),
        (
            "LOD",
            "Derived representation for scalable preview and VR performance evaluation.",
            "Multiple splat budgets or reduced subsets generated from the same source scene.",
            "Used to keep interaction responsive while preserving a comparable scene layout.",
        ),
    ]

    rule(slide, x0, y0, table_w, RULE)
    for x, w, h in zip(col_x, col_w, headers):
        text_box(slide, x, y0 + 0.08, w, 0.18, h, 6.5, True, BLACK)
    rule(slide, x0, y0 + 0.36, table_w, RULE)

    row_y = y0 + 0.48
    row_h = 0.68
    for i, row in enumerate(rows):
        y = row_y + i * row_h
        text_box(slide, col_x[0], y, col_w[0], row_h - 0.08, row[0], 6.9, True, BLACK)
        text_box(slide, col_x[1], y, col_w[1], row_h - 0.08, row[1], 6.25, False, BLACK)
        text_box(slide, col_x[2], y, col_w[2], row_h - 0.08, row[2], 6.25, False, BLACK)
        text_box(slide, col_x[3], y, col_w[3], row_h - 0.08, row[3], 6.25, False, BLACK)
        if i < len(rows) - 1:
            rule(slide, x0, y + row_h - 0.05, table_w, LIGHT_RULE)
    rule(slide, x0, row_y + len(rows) * row_h - 0.05, table_w, RULE)

    text_box(
        slide,
        0.64,
        4.82,
        8.65,
        0.18,
        "Key point: format support is part of the demo contribution because it makes the same nonlinear rendering pipeline usable across captured and exported 3DGS scenes.",
        6.7,
        True,
        TEAL,
        PP_ALIGN.CENTER,
    )
    text_box(
        slide,
        0.64,
        5.06,
        8.65,
        0.12,
        "Note: LOD is treated as a derived runtime representation rather than an external file container.",
        5.4,
        False,
        MUTED,
        PP_ALIGN.CENTER,
    )
    footer(slide)

    prs.save(OUT)
    strip_effects(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
