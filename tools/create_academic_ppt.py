from pathlib import Path
from zipfile import ZipFile, ZIP_DEFLATED
from tempfile import NamedTemporaryFile
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE


ROOT = Path(r"E:\3dgs fisheye projection")
OUT_DIR = ROOT / "Presentation"
OUT_DIR.mkdir(exist_ok=True)
OUT = OUT_DIR / "3DGS_FOV_Fisheye_VR_Academic_Pre.pptx"


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BLACK = RGBColor(18, 18, 18)
GREY = RGBColor(86, 86, 86)
LIGHT = RGBColor(245, 247, 250)
MID = RGBColor(215, 222, 230)
TEAL = RGBColor(0, 135, 145)
BLUE = RGBColor(43, 93, 173)
RED = RGBColor(196, 65, 62)
GREEN = RGBColor(68, 150, 105)
GOLD = RGBColor(196, 151, 53)
PURPLE = RGBColor(113, 82, 165)


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color=RGBColor(255, 255, 255)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, x, y, w, h, text="", font_size=24, bold=False,
                color=BLACK, align=PP_ALIGN.LEFT, font="Aptos"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None, section=None):
    if section:
        add_textbox(slide, 0.72, 0.28, 2.2, 0.25, section.upper(), 8, True, TEAL)
    add_textbox(slide, 0.72, 0.55, 11.8, 0.55, title, 28, True, BLACK)
    if subtitle:
        add_textbox(slide, 0.75, 1.08, 10.8, 0.35, subtitle, 12, False, GREY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.43), Inches(0.78), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def add_footer(slide, n=None):
    add_textbox(slide, 0.72, 7.08, 5.8, 0.18, "3DGS FOV/Fisheye Rendering for VR Scene Exploration", 7, False, RGBColor(130, 130, 130))
    if n is not None:
        add_textbox(slide, 12.45, 7.08, 0.25, 0.18, str(n), 7, False, RGBColor(130, 130, 130), PP_ALIGN.RIGHT)


def bullet_list(slide, x, y, w, h, items, font_size=16, color=BLACK, spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.04)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(5 * spacing)
    return box


def label(slide, x, y, w, h, text, fill=LIGHT, line=MID, color=BLACK,
          font_size=13, bold=False, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    try:
        shape.shadow.inherit = False
    except Exception:
        pass
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return shape


def arrow(slide, x1, y1, x2, y2, color=GREY, width=1.5):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    conn.line.end_arrowhead = True
    return conn


def placeholder(slide, x, y, w, h, title, subtitle=None, accent=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
    shape.line.color.rgb = MID
    shape.line.width = Pt(1.2)
    try:
        shape.shadow.inherit = False
    except Exception:
        pass
    add_textbox(slide, x + 0.18, y + h / 2 - 0.22, w - 0.36, 0.25, title, 13, True, accent, PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, x + 0.18, y + h / 2 + 0.05, w - 0.36, 0.25, subtitle, 9, False, GREY, PP_ALIGN.CENTER)
    return shape


def remove_all_shadows(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                shape.shadow.inherit = False
            except Exception:
                pass


def strip_visual_effects_from_pptx(path: Path):
    """Remove all DrawingML effect nodes so exported shapes stay flat."""
    effect_tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)

    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in effect_tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)

    tmp_path.replace(path)


def pill(slide, x, y, text, color=TEAL, w=1.55):
    return label(slide, x, y, w, 0.33, text, fill=RGBColor(238, 250, 251), line=color, color=color, font_size=10, bold=True)


def add_simple_splat_diagram(slide, x=6.7, y=1.72):
    label(slide, x, y, 1.3, 0.42, "Mesh", fill=RGBColor(248, 248, 248), font_size=12, bold=True)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x + 0.22), Inches(y + 0.65), Inches(0.7), Inches(0.65))
    tri.fill.solid(); tri.fill.fore_color.rgb = RGBColor(232, 238, 247); tri.line.color.rgb = BLUE
    label(slide, x + 1.75, y, 1.3, 0.42, "Point cloud", fill=RGBColor(248, 248, 248), font_size=12, bold=True)
    for dx, dy in [(0.05, 0.2), (0.35, 0.42), (0.72, 0.25), (0.95, 0.55), (0.48, 0.78), (0.18, 0.62)]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.95 + dx), Inches(y + 0.55 + dy), Inches(0.08), Inches(0.08))
        dot.fill.solid(); dot.fill.fore_color.rgb = RED; dot.line.fill.background()
    label(slide, x + 3.55, y, 1.3, 0.42, "Splats", fill=RGBColor(248, 248, 248), font_size=12, bold=True)
    for dx, dy, col in [(0.08, 0.2, TEAL), (0.55, 0.43, GOLD), (0.25, 0.72, PURPLE)]:
        ell = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 3.73 + dx), Inches(y + 0.57 + dy), Inches(0.52), Inches(0.25))
        ell.rotation = 25
        ell.fill.solid(); ell.fill.fore_color.rgb = col; ell.fill.transparency = 25
        ell.line.color.rgb = col


def add_pipeline(slide, y, labels, colors=None):
    x = 0.95
    w = 1.45
    gap = 0.42
    colors = colors or [TEAL] * len(labels)
    for i, text in enumerate(labels):
        label(slide, x + i * (w + gap), y, w, 0.75, text, fill=RGBColor(247, 250, 252), line=colors[i], color=BLACK, font_size=10, bold=True)
        if i < len(labels) - 1:
            arrow(slide, x + w + i * (w + gap) + 0.04, y + 0.37, x + (i + 1) * (w + gap) - 0.08, y + 0.37, GREY, 1.1)


def add_table(slide, x, y, col_widths, row_h, data, header=True, font_size=9):
    rows = len(data)
    cols = len(data[0])
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(sum(col_widths)), Inches(row_h * rows)).table
    for c, width in enumerate(col_widths):
        table.columns[c].width = Inches(width)
    for r in range(rows):
        table.rows[r].height = Inches(row_h)
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(238, 246, 247) if header and r == 0 else RGBColor(255, 255, 255)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(font_size)
                p.font.bold = header and r == 0
                p.font.color.rgb = BLACK
    return table


slides = []


def new_slide(title=None, subtitle=None, section=None):
    slide = blank_slide()
    set_bg(slide)
    if title:
        add_title(slide, title, subtitle, section)
    slides.append(slide)
    return slide


# 1
s = blank_slide(); set_bg(s); slides.append(s)
add_textbox(s, 0.65, 0.75, 11.8, 1.25,
            "FOV and Fisheye Rendering of\n3D Gaussian Splats for VR Scene Exploration",
            34, False, BLACK)
add_textbox(s, 0.68, 2.05, 10.5, 0.4,
            "From asset formats to covariance-aware nonlinear projection in Unity",
            16, False, GREY)
for i, txt in enumerate(["PLY / SPZ / SOG", "Unity + URP", "Desktop + VR", "Cubemap vs Direct"]):
    pill(s, 0.7 + i * 2.15, 2.72, txt, [TEAL, BLUE, GREEN, PURPLE][i], w=1.8)
placeholder(s, 8.55, 3.18, 3.9, 2.3, "Insert final demo screenshot / video still",
            "FOV + fisheye transformation in a Gaussian splat scene", TEAL)
add_textbox(s, 0.7, 5.85, 6.8, 0.35, "Presenter: [Your Name]  ·  Advisor: [Advisor Name]  ·  [Lab / University]", 12, False, BLACK)
add_textbox(s, 0.7, 6.33, 5.5, 0.28, "Group presentation draft · generated structure, editable in PowerPoint", 9, False, GREY)

# 2
s = new_slide("Table of Contents", "The talk is organized as a story: motivation → graphics background → implementation journey → comparison.", "Roadmap")
toc_items = [
    ("1", "Motivation", "Slides 3-4", "Why FOV and fisheye transformations are useful for VR scene exploration."),
    ("2", "Background", "Slides 5-7", "What makes Gaussian splats different from meshes and points."),
    ("3", "Project Setup", "Slides 8-12", "How PLY, SPZ, and SOG assets enter a Unity/VR workflow."),
    ("4", "Implementation Story", "Slides 13-24", "Direct center projection, debugging, cubemap baseline, and covariance-aware direct rendering."),
    ("5", "Comparison", "Slides 25-28", "FOV vs fisheye, cubemap vs direct, desktop vs VR preview."),
    ("6", "Status and Next Steps", "Slides 29-32", "What is done, what remains for ISMAR Demo, and references.")
]
for i, (num, heading, slide_range, desc) in enumerate(toc_items):
    y = 1.62 + i * 0.78
    color = [TEAL, BLUE, PURPLE, GOLD, GREEN, RED][i]
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(y + 0.08), Inches(0.06), Inches(0.48))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_textbox(s, 1.18, y, 0.38, 0.35, num, 16, True, color, PP_ALIGN.RIGHT)
    add_textbox(s, 1.75, y, 3.2, 0.33, heading, 16, True, BLACK)
    add_textbox(s, 5.2, y + 0.03, 1.4, 0.28, slide_range, 9, False, color)
    add_textbox(s, 6.65, y + 0.02, 5.25, 0.33, desc, 10, False, GREY)
add_textbox(s, 1.0, 6.65, 10.8, 0.35,
            "Main narrative: a practical exploration of how nonlinear view transformations become real inside a 3DGS Unity/VR renderer.",
            11, False, TEAL)
add_footer(s, 2)

# 3
s = new_slide("Why Change the View in VR?", "FOV and fisheye transformations can bring peripheral or hidden scene content into view.", "Motivation")
bullet_list(s, 0.85, 1.75, 5.6, 3.2, [
    "VR users have a limited instantaneous field of view.",
    "Some scene content may be peripheral, occluded, or hard to inspect.",
    "Nonlinear viewing can reveal more spatial context without moving the whole scene.",
    "Our target: interactive scene exploration for Gaussian splat environments."
], 17)
placeholder(s, 7.05, 1.65, 2.25, 1.55, "Normal view", "Insert screenshot")
placeholder(s, 9.75, 1.65, 2.25, 1.55, "FOV / fisheye view", "Insert screenshot")
arrow(s, 9.28, 2.42, 9.72, 2.42, TEAL, 2.2)
placeholder(s, 7.05, 4.05, 4.95, 1.45, "Visual message", "Objects near the boundary can be brought into view.")
add_footer(s, 3)

# 4
s = new_slide("Inspirations: Flexible Spatial Transformations", "Advisor demos motivate the kind of transformations we want to support for 3DGS scenes.", "Motivation")
placeholder(s, 0.9, 1.75, 5.55, 3.3, "Insert fisheye demo video / screenshot",
            "https://filonik.pages.lisn.upsaclay.fr/test-pages/demos/fisheye.html", BLUE)
placeholder(s, 6.9, 1.75, 5.55, 3.3, "Insert inception demo video / screenshot",
            "https://filonik.pages.lisn.upsaclay.fr/test-pages/demos/inception.html", PURPLE)
bullet_list(s, 1.2, 5.55, 10.7, 0.8, [
    "Question: can similar viewpoint transformations be applied to photorealistic Gaussian splat scenes in VR?"
], 16, TEAL)
add_footer(s, 4)

# 5
s = new_slide("What Is 3D Gaussian Splatting?", "A 3DGS scene is represented by many anisotropic transparent Gaussians, not a triangle mesh.", "Background")
bullet_list(s, 0.85, 1.7, 5.25, 3.2, [
    "Each splat has position, color, opacity, scale, rotation, and covariance.",
    "Rendering projects each 3D Gaussian into a 2D elliptical footprint.",
    "The image is formed by sorting and alpha-blending many splats.",
    "This makes projection changes more delicate than for meshes or points."
], 16)
add_simple_splat_diagram(s, 6.65, 1.8)
placeholder(s, 6.85, 4.35, 5.0, 1.2, "Insert one cool 3DGS example", "Optional image/video still")
add_footer(s, 5)

# 6
s = new_slide("Why Nonlinear Projection Is Hard for Splats", "Changing only the center position is not enough: the projected footprint also changes.", "Background")
label(s, 0.95, 1.85, 2.3, 0.75, "Mesh\nprojection matrix", fill=RGBColor(248, 248, 248), line=BLUE, font_size=12, bold=True)
label(s, 4.0, 1.85, 2.3, 0.75, "Point cloud\ncenter projection", fill=RGBColor(248, 248, 248), line=RED, font_size=12, bold=True)
label(s, 7.05, 1.85, 2.3, 0.75, "Gaussian splat\ncenter + footprint", fill=RGBColor(248, 248, 248), line=TEAL, font_size=12, bold=True)
arrow(s, 3.35, 2.23, 3.93, 2.23)
arrow(s, 6.42, 2.23, 7.0, 2.23)
placeholder(s, 9.95, 1.55, 2.0, 1.3, "2D ellipse", "projected covariance")
bullet_list(s, 1.1, 3.65, 10.6, 1.5, [
    "For Gaussian splats, a fisheye/FOV change affects the projected 2D covariance.",
    "If the footprint is not transformed consistently, artifacts appear near the image boundary."
], 18)
add_footer(s, 6)

# 7
s = new_slide("3DGS Rendering Pipeline", "The covariance projection stage is the key place where nonlinear viewing becomes visible.", "Background")
add_pipeline(s, 2.1, [
    "3D\nGaussians", "View\ntransform", "2D\ncovariance", "Ellipse\nsplat", "Sort", "Alpha\nblend"
], [TEAL, BLUE, PURPLE, GOLD, RED, GREEN])
placeholder(s, 1.1, 4.3, 10.9, 1.25, "Takeaway", "A nonlinear projection must handle both splat centers and their screen-space ellipses.")
add_footer(s, 7)

# 8
s = new_slide("Project Goal", "Build a Unity/VR prototype for interactive FOV and fisheye transformations of Gaussian splat scenes.", "Our Project")
add_pipeline(s, 1.95, ["Formats", "Unity\nasset", "Renderer", "FOV +\nfisheye", "Desktop\npreview", "VR\npreview"],
             [BLUE, BLUE, TEAL, TEAL, GREEN, GREEN])
bullet_list(s, 1.0, 4.05, 10.8, 1.35, [
    "The project is a system demo: asset pipeline + rendering approaches + interactive preview.",
    "We use the implementation process itself as the story: what broke, why, and how we fixed it."
], 17)
add_footer(s, 8)

# 9
s = new_slide("Challenge 1: Gaussian Splat Formats", "Different tools use different containers and encodings, so loading scenes is part of the research prototype.", "Formats")
add_table(s, 0.85, 1.75, [1.0, 3.0, 3.0, 3.2], 0.55, [
    ["Format", "Where it appears", "Main challenge", "Our status"],
    ["PLY", "Original 3DGS / common exports", "Large files; high memory use", "Importer + max-splats preview"],
    ["SPZ", "Scaniverse / Niantic-style compact splats", "Binary decoding; scale/SH conversion", "Supported through UnityGaussianSplatting"],
    ["SOG", "PlayCanvas / SuperSplat ecosystem", "ZIP + WebP + quantization/codebooks", "Integrated SOG importer; color/memory fixes"],
], True, 10)
placeholder(s, 1.25, 5.35, 10.6, 0.65, "Message", "The demo is not tied to one file format; scenes can move across common 3DGS toolchains.")
add_footer(s, 9)

# 10
s = new_slide("Our Asset Pipeline", "PLY/SPZ/SOG are converted into Unity GaussianSplatAsset data used by the runtime renderer.", "Formats")
add_pipeline(s, 2.0, ["PLY", "SPZ", "SOG\n(WebP/ZIP)", "Importer /\nconverter", "Gaussian\nSplat Asset", "Unity\nRenderer"],
             [BLUE, BLUE, PURPLE, TEAL, TEAL, GREEN])
bullet_list(s, 0.95, 4.1, 11.2, 1.6, [
    "Added SOG package integration and editor import workflow.",
    "Fixed practical issues encountered during conversion: color decoding, memory pressure, and large-scene preview.",
    "Removed generated splat data from version control; repository keeps the reusable pipeline."
], 16)
add_footer(s, 10)

# 11
s = new_slide("Existing Tools and Where We Fit", "Our project sits between web viewers, Unity renderers, and training pipelines.", "Context")
add_table(s, 0.85, 1.7, [2.0, 3.0, 3.0, 3.1], 0.55, [
    ["Tool / line of work", "Strength", "Limitation for us", "Our use"],
    ["UnityGaussianSplatting", "Unity renderer and asset tools", "No built-in VR nonlinear projection workflow", "Base renderer"],
    ["PlayCanvas / SuperSplat", "SOG ecosystem; FOV + fisheye for splats", "Web engine; fisheye disabled in XR in engine source", "Reference method"],
    ["Nerfstudio", "Training and export ecosystem", "Not a Unity/VR runtime demo", "Related pipeline context"],
], True, 9)
add_footer(s, 11)

# 12
s = new_slide("Reference: PlayCanvas FOV + Fisheye", "PlayCanvas provides the closest existing runtime reference for Gaussian splat fisheye.", "Context")
bullet_list(s, 0.85, 1.72, 5.55, 3.4, [
    "Supports camera FOV and gsplat fisheye controls.",
    "Fisheye affects Gaussian splats and infinite sky, not ordinary mesh/UI.",
    "Recommends radial sorting for fisheye / very wide views.",
    "Uses a generalized fisheye model and covariance-aware splat rendering."
], 16)
placeholder(s, 7.0, 1.72, 4.9, 2.2, "Insert PlayCanvas screenshot / formula", "g(theta) = k · tan(theta / k)")
placeholder(s, 7.0, 4.35, 4.9, 1.25, "Our question", "How does this translate to Unity + VR, and how does it compare to cubemap distortion?")
add_footer(s, 12)

# 13
s = new_slide("Attempt 1: Direct Center Fisheye", "First working version: transform splat centers with a fisheye mapping.", "Implementation Story")
placeholder(s, 0.95, 1.75, 5.3, 3.15, "Insert first direct fisheye result", "Worked visually in the center")
placeholder(s, 6.85, 1.75, 5.3, 3.15, "Insert edge artifact example", "Problems near image boundary")
bullet_list(s, 1.15, 5.35, 10.6, 0.8, [
    "This was useful because it isolated the first problem: center projection alone does not make splats correct."
], 16, TEAL)
add_footer(s, 13)

# 14
s = new_slide("Problem: Splats at the Edge", "The center can be correct while the ellipse footprint is wrong.", "Implementation Story")
label(s, 1.0, 2.0, 2.2, 0.75, "Center\nprojection", line=GREEN, fill=RGBColor(247, 252, 249), font_size=13, bold=True)
label(s, 4.0, 2.0, 2.2, 0.75, "Projected\ncovariance", line=RED, fill=RGBColor(253, 247, 247), font_size=13, bold=True)
label(s, 7.0, 2.0, 2.2, 0.75, "Culling /\nsorting", line=GOLD, fill=RGBColor(253, 250, 242), font_size=13, bold=True)
arrow(s, 3.25, 2.37, 3.95, 2.37); arrow(s, 6.25, 2.37, 6.95, 2.37)
bullet_list(s, 1.05, 3.65, 10.7, 1.45, [
    "Artifacts suggested that the projection of the Gaussian footprint had not been adapted.",
    "Wide-angle/fisheye views also break the usual assumption that camera Z is a globally good transparency order."
], 17)
add_footer(s, 14)

# 15
s = new_slide("Debugging View: Render Splats as Points", "A point view helped verify that splat centers were transformed correctly.", "Implementation Story")
placeholder(s, 0.95, 1.75, 5.25, 3.4, "Insert point-debug view", "Splat centers only")
placeholder(s, 6.85, 1.75, 5.25, 3.4, "Insert splat view with artifacts", "Centers okay; footprints problematic")
bullet_list(s, 1.05, 5.55, 10.8, 0.6, ["Lesson: the failure was not simply the center transform."], 18, TEAL)
add_footer(s, 15)

# 16
s = new_slide("Lesson: A Splat Is Not a Point", "Correct nonlinear rendering requires center projection, footprint projection, sorting, and culling to agree.", "Implementation Story")
add_pipeline(s, 2.05, ["Center", "Covariance", "Sort", "Cull", "Composite"], [GREEN, TEAL, GOLD, RED, BLUE])
placeholder(s, 1.1, 4.25, 10.8, 1.1, "Design implication", "We need either a robust image-space baseline or a renderer-level solution.")
add_footer(s, 16)

# 17
s = new_slide("Attempt 2: Six-Camera Cubemap Fisheye", "Render the scene into perspective faces, then apply fisheye distortion during composition.", "Cubemap Method")
label(s, 0.95, 1.85, 1.2, 0.65, "+X", line=BLUE, font_size=15, bold=True)
label(s, 2.25, 1.85, 1.2, 0.65, "-X", line=BLUE, font_size=15, bold=True)
label(s, 3.55, 1.85, 1.2, 0.65, "+Y", line=BLUE, font_size=15, bold=True)
label(s, 4.85, 1.85, 1.2, 0.65, "-Y", line=BLUE, font_size=15, bold=True)
label(s, 6.15, 1.85, 1.2, 0.65, "+Z", line=BLUE, font_size=15, bold=True)
label(s, 7.45, 1.85, 1.2, 0.65, "-Z", line=BLUE, font_size=15, bold=True)
arrow(s, 8.75, 2.17, 9.6, 2.17, TEAL, 2)
placeholder(s, 9.75, 1.35, 2.55, 1.65, "Composite shader", "fisheye sample direction")
bullet_list(s, 1.05, 4.0, 10.7, 1.35, [
    "This approach is easy to reason about because each capture is ordinary perspective.",
    "It also provides a high-quality baseline for side-by-side comparison."
], 17)
add_footer(s, 17)

# 18
s = new_slide("Desktop Cubemap Implementation", "DesktopHighQualityFisheye renders six faces and composites them to the output view.", "Cubemap Method")
add_pipeline(s, 2.0, ["Output\ncamera", "6 capture\ncameras", "Render\ntextures", "Composite\nshader", "Screen\noutput"],
             [TEAL, BLUE, BLUE, PURPLE, GREEN])
placeholder(s, 1.0, 4.1, 5.2, 1.5, "Insert desktop cubemap demo", "FOV / fisheye controls")
placeholder(s, 6.85, 4.1, 5.2, 1.5, "Insert Unity hierarchy / inspector", "Optional implementation screenshot")
add_footer(s, 18)

# 19
s = new_slide("VR Cubemap Implementation", "VR needs separate cubemap captures for left and right eyes.", "Cubemap Method")
label(s, 1.0, 1.85, 2.0, 0.6, "Left eye", line=TEAL, fill=RGBColor(238, 250, 251), font_size=14, bold=True)
label(s, 1.0, 3.1, 2.0, 0.6, "Right eye", line=PURPLE, fill=RGBColor(247, 244, 252), font_size=14, bold=True)
for i in range(6):
    label(s, 3.4 + i * 0.92, 1.82, 0.62, 0.48, f"F{i+1}", line=TEAL, font_size=9)
    label(s, 3.4 + i * 0.92, 3.07, 0.62, 0.48, f"F{i+1}", line=PURPLE, font_size=9)
arrow(s, 8.95, 2.4, 10.0, 2.4, TEAL, 2)
placeholder(s, 10.15, 1.73, 1.75, 1.45, "Stereo\ncomposite", "12 faces")
bullet_list(s, 1.05, 4.55, 10.8, 1.1, [
    "Implemented double buffering and delayed swapping so incomplete face updates do not flash as black frames.",
    "This made the method usable for VR preview but exposed its cost and latency."
], 15)
add_footer(s, 19)

# 20
s = new_slide("Cubemap Method: Pros and Cons", "Cubemap fisheye is a robust baseline, but expensive in VR.", "Cubemap Method")
add_table(s, 1.05, 1.8, [2.4, 4.2, 4.2], 0.62, [
    ["Aspect", "Strength", "Limitation"],
    ["Visual quality", "Stable at image edges", "Limited by face resolution and seams"],
    ["Implementation", "Uses ordinary perspective captures", "Needs capture camera management"],
    ["VR", "Conceptually simple stereo", "12 faces for two eyes; latency/update complexity"],
    ["Research value", "Good baseline for comparison", "Not the most direct splat rendering solution"],
], True, 10)
add_footer(s, 20)

# 21
s = new_slide("Attempt 3: Direct Covariance-Aware Fisheye", "Final direction: modify the splat renderer itself instead of post-distorting captures.", "Direct Method")
add_pipeline(s, 1.95, ["FOV +\nfisheye params", "Center\nprojection", "Radial\nsorting", "Covariance\nJacobian", "Cull +\nAA", "Composite"],
             [TEAL, GREEN, GOLD, PURPLE, RED, BLUE])
bullet_list(s, 1.0, 4.2, 11.0, 1.35, [
    "Implemented PlayCanvas-style fisheye parameter mapping in Unity.",
    "Projected splat centers through the fisheye model.",
    "Transformed projected covariance using a fisheye-aware Jacobian.",
    "Adjusted sorting/culling behavior for wide-angle rendering."
], 15)
add_footer(s, 21)

# 22
s = new_slide("Projection Model: FOV and Fisheye Are Separate Controls", "FOV changes angular coverage; fisheye changes how angles are compressed onto the screen.", "Direct Method")
add_table(s, 1.0, 1.78, [2.2, 4.0, 4.0], 0.6, [
    ["Control", "Effect", "What the user sees"],
    ["FOV", "Expands or narrows the virtual field of view", "More or less angular scene coverage"],
    ["Fisheye", "Nonlinearly compresses rays toward screen center", "Wide/tiny-planet-like distortion"],
    ["FOV + fisheye", "Combines coverage and nonlinear compression", "Peripheral content becomes visible"],
], True, 10)
placeholder(s, 1.0, 4.55, 2.55, 1.25, "Normal", "Insert screenshot")
placeholder(s, 3.85, 4.55, 2.55, 1.25, "High FOV", "Insert screenshot")
placeholder(s, 6.7, 4.55, 2.55, 1.25, "Fisheye", "Insert screenshot")
placeholder(s, 9.55, 4.55, 2.55, 1.25, "FOV + fisheye", "Insert screenshot")
add_footer(s, 22)

# 23
s = new_slide("Covariance Transform", "The 3D Gaussian footprint must be projected through the nonlinear mapping.", "Direct Method")
label(s, 1.0, 2.1, 2.2, 0.75, "3D Gaussian\ncovariance", line=TEAL, fill=RGBColor(238, 250, 251), font_size=13, bold=True)
arrow(s, 3.35, 2.48, 4.35, 2.48, GREY, 1.6)
label(s, 4.55, 2.1, 2.2, 0.75, "Fisheye\nJacobian", line=PURPLE, fill=RGBColor(247, 244, 252), font_size=13, bold=True)
arrow(s, 6.9, 2.48, 7.9, 2.48, GREY, 1.6)
label(s, 8.1, 2.1, 2.2, 0.75, "2D ellipse\nfootprint", line=GOLD, fill=RGBColor(253, 250, 242), font_size=13, bold=True)
add_textbox(s, 1.15, 3.75, 10.8, 0.85,
            "This is the graphics principle behind the final direct approach: the image-space ellipse is no longer the same as in perspective projection.",
            17, False, BLACK, PP_ALIGN.CENTER)
add_footer(s, 23)

# 24
s = new_slide("Sorting, Culling, and VR Comfort", "Wide-angle splat rendering changes practical renderer assumptions.", "Direct Method")
bullet_list(s, 0.9, 1.75, 5.55, 3.5, [
    "Depth sorting works well for a narrow forward-facing perspective view.",
    "In fisheye views, side and near-surround content may be visible.",
    "Radial distance sorting is more stable for panoramic/fisheye images.",
    "Near fade reduces close-to-eye splat discomfort in VR."
], 15)
placeholder(s, 7.05, 1.85, 4.75, 1.3, "Depth sorting", "good for narrow perspective")
placeholder(s, 7.05, 3.65, 4.75, 1.3, "Radial sorting", "more stable for wide/fisheye")
add_footer(s, 24)

# 25
s = new_slide("Desktop and VR Preview System", "The prototype includes runtime controls and preview scenes for both desktop and VR.", "System")
add_table(s, 0.95, 1.75, [2.5, 4.0, 4.1], 0.58, [
    ["Component", "Purpose", "Status"],
    ["DesktopPreview", "Fast iteration without headset", "Implemented"],
    ["XRSimulatorTemplate", "Editor VR-style movement and preview", "Implemented"],
    ["Projection control panel", "Sliders for FOV and fisheye", "Implemented"],
    ["Keyboard controls", "Quick FOV/fisheye reset and adjustment", "Implemented"],
    ["VR high-quality fisheye", "Stereo cubemap preview", "Implemented"],
], True, 10)
add_footer(s, 25)

# 26
s = new_slide("Comparison Plan", "We compare user-visible behavior and renderer trade-offs, not only final screenshots.", "Evaluation")
add_table(s, 0.55, 1.65, [1.55, 1.75, 1.6, 1.8, 1.8, 2.1], 0.52, [
    ["Mode", "Coverage", "Distortion", "Edge quality", "VR cost", "Notes"],
    ["Perspective", "low", "none", "stable", "low", "baseline"],
    ["High FOV", "medium/high", "low", "can stretch", "low", "simple control"],
    ["Cubemap fisheye", "high", "high", "stable", "high", "6/12 captures"],
    ["Direct fisheye", "high", "high", "depends on covariance/sorting", "lower potential", "renderer-level"],
], True, 8.5)
placeholder(s, 1.0, 5.15, 10.8, 0.65, "Insert measured numbers later", "FPS / GPU time / resolution / scene size / headset or desktop hardware")
add_footer(s, 26)

# 27
s = new_slide("Side-by-Side Results Placeholder", "This is the key slide for the final presentation: same scene, same viewpoint, different rendering modes.", "Results")
placeholder(s, 0.75, 1.72, 2.9, 1.75, "Normal perspective", "Insert screenshot")
placeholder(s, 3.95, 1.72, 2.9, 1.75, "High FOV", "Insert screenshot")
placeholder(s, 7.15, 1.72, 2.9, 1.75, "Cubemap fisheye", "Insert screenshot")
placeholder(s, 10.35, 1.72, 2.25, 1.75, "Direct fisheye", "Insert screenshot")
placeholder(s, 1.05, 4.45, 5.1, 1.35, "Insert demo video: cubemap path", "Desktop or VR")
placeholder(s, 7.15, 4.45, 5.1, 1.35, "Insert demo video: direct path", "Desktop or VR")
add_footer(s, 27)

# 28
s = new_slide("Comparison With Other Implementations", "Our contribution is a Unity/VR system and implementation comparison, not claiming first fisheye splat rendering.", "Context")
add_table(s, 0.65, 1.58, [2.0, 2.65, 2.5, 3.4, 1.65], 0.5, [
    ["System", "FOV/fisheye", "VR/XR", "Format/runtime context", "Role"],
    ["PlayCanvas", "Yes, splats + sky", "Fisheye disabled in XR source", "Web/SOG ecosystem", "Reference"],
    ["krpano", "Fisheye/little planet with GS support", "Panorama engine context", "Tour/panorama viewer", "Related demo"],
    ["UnityGaussianSplatting", "No dedicated fisheye workflow", "Unity runtime", "PLY/SPZ assets", "Base"],
    ["Our prototype", "FOV + fisheye; cubemap + direct", "Desktop + VR preview", "PLY/SPZ/SOG in Unity", "Demo"],
], True, 8.5)
add_footer(s, 28)

# 29
s = new_slide("Current Status", "What is already implemented in the repository.", "Status")
items = [
    "Unity project with embedded Gaussian Splatting packages",
    "PLY/SPZ and SOG asset import pipeline",
    "SOG color decoding and memory/import fixes",
    "Desktop preview scene and camera controls",
    "XR simulator scene and projection control panel",
    "Cubemap fisheye path for desktop and stereo VR",
    "Direct PlayCanvas-style fisheye projection branch",
    "Fisheye-aware covariance transform and radial sorting",
    "Placeholders for side-by-side comparison and demo video"
]
bullet_list(s, 0.95, 1.65, 11.2, 4.7, [f"✓ {i}" for i in items], 13)
add_footer(s, 29)

# 30
s = new_slide("Next Steps Toward ISMAR Demo", "The remaining work is mostly evidence and packaging.", "Status")
add_table(s, 1.0, 1.75, [3.0, 5.4, 2.2], 0.58, [
    ["Task", "Output", "Priority"],
    ["Record videos", "Short clips for FOV, fisheye, cubemap, direct, VR", "High"],
    ["Capture comparison figures", "Same viewpoint across modes", "High"],
    ["Measure performance", "FPS/GPU time for scene sizes and resolutions", "High"],
    ["Write demo abstract", "100-word abstract + 2-page extended abstract", "High"],
    ["Prepare teaser", "Landscape image with four-mode comparison", "Medium"],
], True, 10)
add_footer(s, 30)

# 31
s = new_slide("Takeaways", "What people should learn from the project.", "Conclusion")
bullet_list(s, 1.0, 1.75, 10.8, 3.8, [
    "Format conversion is part of the practical 3DGS workflow, not just a setup detail.",
    "Changing a 3DGS view is harder than changing a camera matrix because splat footprints matter.",
    "A cubemap approach is robust and easy to validate, but costly in VR.",
    "A direct covariance-aware approach is more renderer-native and a better long-term path.",
    "FOV and fisheye transformations can support immersive scene exploration by bringing peripheral content into view."
], 18)
add_footer(s, 31)

# 32
s = new_slide("References", "Keep this slide at the end; cite only the most important sources in the spoken presentation.", "References")
refs = [
    "Kerbl et al. 3D Gaussian Splatting for Real-Time Radiance Field Rendering. ACM TOG / SIGGRAPH, 2023.",
    "Aras Pranckevičius. UnityGaussianSplatting. GitHub: github.com/aras-p/UnityGaussianSplatting.",
    "PlayCanvas. Fisheye Rendering for Gaussian Splats. developer.playcanvas.com/user-manual/gaussian-splatting/building/fisheye/.",
    "PlayCanvas. SOG Format. developer.playcanvas.com/user-manual/gaussian-splatting/formats/sog/.",
    "PlayCanvas Engine source. FisheyeProjection and gsplat shaders. github.com/playcanvas/engine.",
    "Tancik et al. Nerfstudio: A Modular Framework for Neural Radiance Field Development. SIGGRAPH, 2023.",
    "Wu et al. 3DGUT: Enabling Distorted Cameras and Secondary Rays in Gaussian Splatting. CVPR, 2025.",
    "Yang et al. Fisheye-GS: Lightweight and Extensible Gaussian Splatting Module for Fisheye Cameras. arXiv, 2024.",
    "krpano 1.23 3D Gaussian Splatting support. krpano.com/forum."
]
bullet_list(s, 0.85, 1.55, 11.9, 5.3, refs, 10)
add_footer(s, 32)

# Ensure footer numbering on title omitted and all standard slides already handled.

remove_all_shadows(prs)
prs.save(OUT)
strip_visual_effects_from_pptx(OUT)
print(OUT)
