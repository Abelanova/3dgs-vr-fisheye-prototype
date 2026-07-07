from pathlib import Path
from zipfile import ZipFile, ZIP_DEFLATED
from tempfile import NamedTemporaryFile

from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt


SRC = Path(r"C:\Users\nova\Downloads\FOV and Fisheye Distortion Rendering of 3D Gaussian Splats for VR Scene Exploration (1).pptx")
OUT = Path(r"E:\3dgs fisheye projection\Presentation\FOV_Fisheye_3DGS_with_LaTeX_formulas.pptx")
FORMULA_DIR = OUT.parent / "formula_renders"

BLACK = RGBColor(20, 20, 20)
GREY = RGBColor(85, 85, 85)
TEAL = RGBColor(0, 135, 145)
BLUE = RGBColor(43, 93, 173)
PURPLE = RGBColor(113, 82, 165)
LIGHT_TEAL = RGBColor(237, 250, 251)
LIGHT_BLUE = RGBColor(238, 244, 252)
LIGHT_PURPLE = RGBColor(246, 242, 251)


def render_formula(expr: str, out: Path, fontsize=26, color="#141414"):
    """Render a LaTeX-style math expression to a transparent PNG."""
    out.parent.mkdir(parents=True, exist_ok=True)
    fig = plt.figure(figsize=(0.01, 0.01), dpi=300)
    text = fig.text(0, 0, expr, fontsize=fontsize, color=color)
    fig.canvas.draw()
    bbox = text.get_window_extent()
    w = max(bbox.width / fig.dpi + 0.18, 0.5)
    h = max(bbox.height / fig.dpi + 0.16, 0.25)
    plt.close(fig)

    fig = plt.figure(figsize=(w, h), dpi=300)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")
    ax.text(0.5, 0.5, expr, ha="center", va="center", fontsize=fontsize, color=color)
    fig.savefig(out, transparent=True, bbox_inches="tight", pad_inches=0.03)
    plt.close(fig)


def add_text(slide, x, y, w, h, text, size=9, bold=False, color=GREY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_latex_box(slide, x, y, w, h, title, image_path, note=None, accent=TEAL, fill=LIGHT_TEAL, img_h=0.36):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = accent
    box.line.width = Pt(1.1)
    try:
        box.shadow.inherit = False
    except Exception:
        pass

    add_text(slide, x + 0.12, y + 0.07, w - 0.24, 0.18, title, 8.5, True, accent)
    slide.shapes.add_picture(str(image_path), Inches(x + 0.22), Inches(y + 0.31), width=Inches(w - 0.44), height=Inches(img_h))
    if note:
        add_text(slide, x + 0.14, y + h - 0.27, w - 0.28, 0.18, note, 7.2, False, GREY)
    return box


def strip_visual_effects_from_pptx(path: Path):
    effect_tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)
    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in effect_tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(path)


FORMULAS = {
    "gaussian": r"$G(\mathbf{x})=\exp\!\left[-\frac{1}{2}(\mathbf{x}-\boldsymbol{\mu})^{T}\boldsymbol{\Sigma}^{-1}(\mathbf{x}-\boldsymbol{\mu})\right]$",
    "fov": r"$f_y=\frac{1}{\tan(\mathrm{FOV}_y/2)}$",
    "fish_theta": r"$r=g(\theta),\quad \theta=\operatorname{atan2}\!\left(\sqrt{x^2+y^2},-z\right)$",
    "cov": r"$\boldsymbol{\Sigma}_{2D}=J\,W\,\boldsymbol{\Sigma}_{3D}\,W^{T}J^{T}$",
    "fish": r"$g(\theta)=k\,\tan\!\left(\frac{\theta}{k}\right)$",
    "center_fail": r"$\boldsymbol{\mu}\mapsto\pi_{\mathrm{fish}}(\boldsymbol{\mu})\;\not\Rightarrow\;\boldsymbol{\Sigma}_{2D}\ \mathrm{is\ correct}$",
    "direct": r"$\boldsymbol{\Sigma}_{2D}^{\mathrm{fish}}=J_{\mathrm{fish}}\,W\,\boldsymbol{\Sigma}_{3D}\,W^{T}J_{\mathrm{fish}}^{T}$",
    "cubemap": r"$\mathrm{color}(u,v)=C_{\mathrm{face}}\!\left(\mathbf{d}(u,v)\right)$",
    "sorting": r"$\mathrm{Perspective:}\ z_{\mathrm{view}}\qquad\quad\mathrm{Fisheye:}\ \|\mathbf{p}_{\mathrm{view}}\|$",
}


def main():
    if not SRC.exists():
        raise FileNotFoundError(SRC)
    OUT.parent.mkdir(exist_ok=True)
    FORMULA_DIR.mkdir(exist_ok=True)

    rendered = {}
    for name, expr in FORMULAS.items():
        p = FORMULA_DIR / f"{name}.png"
        render_formula(expr, p, fontsize=28 if name not in {"sorting"} else 24)
        rendered[name] = p

    prs = Presentation(SRC)
    slides = prs.slides

    if len(slides) >= 3:
        add_latex_box(slides[2], 7.05, 4.62, 4.95, 1.18, "3D Gaussian representation",
                      rendered["gaussian"], "μ: splat center,  Σ: anisotropic 3D shape", TEAL, LIGHT_TEAL, img_h=0.42)

    if len(slides) >= 8:
        add_latex_box(slides[7], 0.92, 5.28, 5.35, 1.2, "Rectilinear FOV",
                      rendered["fov"], "FOV near 180° causes extreme boundary stretching.", BLUE, LIGHT_BLUE, img_h=0.42)
        add_latex_box(slides[7], 6.92, 5.28, 5.35, 1.2, "Fisheye radial mapping",
                      rendered["fish_theta"], "Fisheye maps viewing angle instead of a flat image plane.", PURPLE, LIGHT_PURPLE, img_h=0.42)

    if len(slides) >= 9:
        add_latex_box(slides[8], 1.15, 5.12, 10.65, 1.02, "Projected covariance",
                      rendered["cov"], "W is the view transform; J is the projection Jacobian.", TEAL, LIGHT_TEAL, img_h=0.38)

    if len(slides) >= 15:
        add_latex_box(slides[14], 7.02, 3.00, 4.88, 1.12, "Generalized fisheye model",
                      rendered["fish"], "Reference model used by PlayCanvas-style fisheye projection.", PURPLE, LIGHT_PURPLE, img_h=0.42)

    if len(slides) >= 17:
        add_latex_box(slides[16], 0.98, 5.28, 10.95, 1.02, "Why center-only projection fails",
                      rendered["center_fail"], "The center can be correct while the ellipse footprint is wrong.", TEAL, LIGHT_TEAL, img_h=0.38)

    if len(slides) >= 20:
        add_latex_box(slides[19], 8.75, 4.92, 3.25, 1.0, "Cubemap sampling",
                      rendered["cubemap"], "Distortion happens after ordinary perspective captures.", BLUE, LIGHT_BLUE, img_h=0.38)

    if len(slides) >= 24:
        add_latex_box(slides[23], 0.98, 5.02, 10.95, 1.02, "Direct fisheye covariance",
                      rendered["direct"], "The direct path changes the splat renderer, not a final image.", TEAL, LIGHT_TEAL, img_h=0.38)

    if len(slides) >= 25:
        add_latex_box(slides[24], 1.0, 5.88, 5.25, 0.82, "FOV controls scale",
                      rendered["fov"], None, BLUE, LIGHT_BLUE, img_h=0.36)
        add_latex_box(slides[24], 6.95, 5.88, 5.25, 0.82, "Fisheye controls radial mapping",
                      rendered["fish"], None, PURPLE, LIGHT_PURPLE, img_h=0.36)

    if len(slides) >= 26:
        add_latex_box(slides[25], 1.35, 5.06, 10.35, 1.02, "Covariance transform",
                      rendered["direct"], "This is the theoretical reason final direct rendering differs from center-only fisheye.", TEAL, LIGHT_TEAL, img_h=0.38)

    if len(slides) >= 27:
        add_latex_box(slides[26], 7.05, 5.00, 4.75, 1.08, "Sorting heuristic",
                      rendered["sorting"], "Radial sorting is more stable when the visible view is very wide.", PURPLE, LIGHT_PURPLE, img_h=0.38)

    prs.save(OUT)
    strip_visual_effects_from_pptx(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
