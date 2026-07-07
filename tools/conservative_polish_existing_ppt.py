from pathlib import Path
from zipfile import ZipFile, ZIP_DEFLATED
from tempfile import NamedTemporaryFile
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor


ROOT = Path(r"E:\3dgs fisheye projection")
SRC = ROOT / "Presentation" / "FOV_Fisheye_3DGS_clean_LaTeX_formulas_aspect.pptx"
OUT = ROOT / "Presentation" / "FOV_Fisheye_3DGS_conservative_polish.pptx"

BLACK = RGBColor(18, 18, 18)
GREY = RGBColor(82, 82, 82)
MUTED = RGBColor(120, 120, 120)
TEAL = RGBColor(0, 135, 145)
BLUE = RGBColor(43, 93, 173)
RED = RGBColor(196, 65, 62)
PURPLE = RGBColor(113, 82, 165)
GREEN = RGBColor(64, 145, 98)
GOLD = RGBColor(186, 142, 45)
MID = RGBColor(215, 222, 230)
PALE = RGBColor(250, 251, 252)


def remove_nonpictures(slide):
    """Preserve every existing picture; remove text/shape scaffolding."""
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        if not hasattr(shape, "image"):
            sp_tree.remove(shape._element)


def clear_slide(slide):
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        sp_tree.remove(shape._element)


def text_box(slide, x, y, w, h, text, size=13, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def bullet_box(slide, x, y, w, h, items, size=12.5, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return box


def title(slide, text, subtitle=None, section=None):
    if section:
        text_box(slide, 0.72, 0.28, 2.4, 0.22, section.upper(), 8, True, TEAL)
    text_box(slide, 0.72, 0.55, 11.7, 0.45, text, 24, True, BLACK)
    if subtitle:
        text_box(slide, 0.74, 1.02, 11.3, 0.28, subtitle, 10.5, False, GREY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.34), Inches(0.7), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def footer(slide, n):
    text_box(slide, 0.72, 7.08, 5.6, 0.18, "3DGS FOV/Fisheye Rendering for VR Scene Exploration", 6.5, False, MUTED)
    text_box(slide, 12.28, 7.08, 0.28, 0.18, str(n), 7, False, MUTED, PP_ALIGN.RIGHT)


def source(slide, text):
    text_box(slide, 0.72, 6.78, 11.7, 0.18, text, 6.5, False, MUTED)


def label(slide, x, y, w, h, text, line=TEAL, fill=RGBColor(255, 255, 255), size=10.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(1)
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = BLACK
    return sh


def arrow(slide, x1, y1, x2, y2, color=GREY):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.2)
    conn.line.end_arrowhead = True
    return conn


def picture_shapes(slide):
    return [s for s in slide.shapes if hasattr(s, "image")]


def set_pic(shape, x, y, w=None, h=None):
    shape.left = Inches(x)
    shape.top = Inches(y)
    if w is not None:
        shape.width = Inches(w)
    if h is not None:
        shape.height = Inches(h)


def strip_effects(path: Path):
    tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)
    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(path)


def polish():
    prs = Presentation(SRC)

    # 1. Title: preserve all pictures; only replace text layout.
    s = prs.slides[0]
    remove_nonpictures(s)
    pics = picture_shapes(s)
    # Keep QR/logo/headshots but push them into tidy places if they exist.
    if len(pics) >= 1:
        # Small logos/images near bottom, QR at right if present.
        for i, p in enumerate(pics):
            if i < 4:
                set_pic(p, 0.82 + i * 1.55, 4.95, w=1.15)
            else:
                set_pic(p, 10.65, 4.75, w=1.25)
    text_box(s, 0.72, 0.72, 11.5, 1.05,
             "FOV and Fisheye Distortion Rendering\nof 3D Gaussian Splats for VR Scene Exploration",
             30, True)
    text_box(s, 0.74, 2.02, 8.9, 0.28,
             "A Unity prototype for nonlinear view transformations of Gaussian splat scenes",
             14, False, GREY)
    for i, (txt, col) in enumerate([("PLY / SPZ / SOG", BLUE), ("Unity + URP", TEAL), ("Desktop + VR", GREEN), ("Cubemap vs Direct", PURPLE)]):
        label(s, 0.74 + i * 2.1, 2.68, 1.7, 0.34, txt, col, size=9.2)
    text_box(s, 0.74, 5.88, 5.6, 0.28, "Yingfangzhong SUN  ·  Supervisor: Daniel Filonik", 12, False)

    # 2. TOC: no pictures, clean short text.
    s = prs.slides[1]
    clear_slide(s)
    title(s, "Talk Roadmap", "From 3DGS representation to nonlinear rendering in Unity/VR.")
    toc = [
        ("1", "Background", "3DGS representation and splat footprints."),
        ("2", "Motivation", "Why FOV/fisheye helps VR scene exploration."),
        ("3", "Formats", "PLY, SPZ, and SOG in a Unity asset pipeline."),
        ("4", "Implementation", "Direct attempt, point debug, cubemap, direct covariance."),
        ("5", "Comparison", "FOV vs fisheye; cubemap vs direct; desktop vs VR."),
        ("6", "Discussion", "Current status, next steps, and references."),
    ]
    for i, (num, head, desc) in enumerate(toc):
        y = 1.7 + i * 0.72
        col = [BLUE, TEAL, PURPLE, GOLD, GREEN, RED][i]
        text_box(s, 0.95, y, 0.32, 0.25, num, 14, True, col, PP_ALIGN.RIGHT)
        text_box(s, 1.55, y, 2.1, 0.25, head, 13.5, True)
        text_box(s, 3.9, y + 0.02, 7.5, 0.25, desc, 10.3, False, GREY)
    footer(s, 2)

    # 3. Background: preserve existing diagram-ish content by only improving text boundaries.
    s = prs.slides[2]
    # Leave shapes and pictures as they are; replace main text boxes only where safe.
    for sh in list(s.shapes):
        if hasattr(sh, "text") and sh.text.strip() and ("Core idea" in sh.text or "3DGS is" in sh.text or "input:" in sh.text or "References:" in sh.text):
            sh.text_frame.clear()
    text_box(s, 0.82, 1.38, 2.8, 0.22, "Core idea", 13, True, TEAL)
    bullet_box(s, 0.82, 1.8, 2.9, 1.7, [
        "Input: posed multi-view images.",
        "Representation: many anisotropic 3D Gaussians.",
        "Rendering: project, sort, and alpha-blend splats."
    ], 11.5)
    text_box(s, 0.82, 3.72, 2.85, 0.42, "A splat has position, opacity, color, and shape.", 10.5, True, TEAL, PP_ALIGN.CENTER)
    source(s, "References: Kerbl et al., 3D Gaussian Splatting, SIGGRAPH 2023; Mildenhall et al., NeRF, ECCV 2020.")

    # 4. Preserve big image, add clean caption/source.
    s = prs.slides[3]
    remove_nonpictures(s)
    pics = picture_shapes(s)
    if pics:
        set_pic(pics[0], 0.92, 1.45, w=5.65)
    title(s, "From Captured Images to a Renderable Scene",
          "3DGS learns an explicit splat representation that can be rendered from new viewpoints.",
          "Background")
    bullet_box(s, 7.0, 1.65, 4.75, 1.7, [
        "Training reconstructs scene radiance from posed images.",
        "Runtime rendering is explicit and interactive.",
        "This makes 3DGS attractive for desktop and VR scene inspection."
    ], 13)
    text_box(s, 7.0, 4.15, 4.8, 0.45, "But changing the view model is harder than changing a mesh camera.", 15, True, TEAL)
    source(s, "Image source: LearnOpenCV 3D Gaussian Splatting explanation. Method reference: Kerbl et al., 2023.")
    footer(s, 4)

    # 5. Preserve screenshot, arrange text to right.
    s = prs.slides[4]
    remove_nonpictures(s)
    pics = picture_shapes(s)
    if pics:
        set_pic(pics[0], 0.85, 1.48, w=6.45)
    title(s, "Why 3DGS Is Interesting for VR",
          "Photorealistic captured spaces can be inspected interactively, but the user still has limited instantaneous view.",
          "Background")
    bullet_box(s, 7.65, 1.65, 4.35, 2.2, [
        "Real-time novel-view rendering.",
        "Scene detail comes from captured images.",
        "Useful for immersive inspection of real spaces.",
        "Question: how should wide FOV and fisheye views work for splats?"
    ], 12.8)
    source(s, "Image source: 3DGS web viewer screenshot from current deck. Verify original URL before final submission.")
    footer(s, 5)

    # 6. Clean motivation, preserve none/placeholder images if present.
    s = prs.slides[5]
    remove_nonpictures(s)
    title(s, "Motivation: Bring More Scene Content Into View",
          "FOV and fisheye controls can reveal peripheral or hidden content during VR scene exploration.",
          "Motivation")
    bullet_box(s, 0.9, 1.75, 4.9, 2.1, [
        "VR has a limited instantaneous field of view.",
        "Relevant content may lie outside the current view.",
        "Nonlinear viewing can compress more spatial context into the display."
    ], 13.5)
    label(s, 6.5, 1.85, 2.1, 0.62, "normal view", BLUE)
    arrow(s, 8.72, 2.16, 9.35, 2.16, TEAL)
    label(s, 9.55, 1.85, 2.1, 0.62, "wide / fisheye view", TEAL)
    text_box(s, 6.4, 3.35, 5.35, 0.55, "Goal: interactive nonlinear viewing transformations for Gaussian splat scenes in Unity and VR.", 15, True, TEAL, PP_ALIGN.CENTER)
    footer(s, 6)

    # 7. Preserve two images, clean captions.
    s = prs.slides[6]
    remove_nonpictures(s)
    pics = picture_shapes(s)
    if len(pics) >= 2:
        set_pic(pics[0], 0.85, 1.48, w=5.35)
        set_pic(pics[1], 6.78, 1.48, w=5.35)
    title(s, "Inspiration: Flexible Spatial Transformations",
          "We first studied simpler demos, then asked how similar transformations behave for Gaussian splats.",
          "Motivation")
    text_box(s, 1.0, 4.65, 4.95, 0.24, "World bending / spatial deformation demo", 10.5, True, BLACK, PP_ALIGN.CENTER)
    text_box(s, 6.95, 4.65, 4.95, 0.24, "FOV / fisheye projection demo", 10.5, True, BLACK, PP_ALIGN.CENTER)
    text_box(s, 1.1, 5.42, 10.8, 0.4, "Challenge for 3DGS: splats are not ordinary geometry, so center positions and screen-space footprints must transform consistently.", 13.2, True, TEAL, PP_ALIGN.CENTER)
    source(s, "Demo sources: advisor-provided fisheye and inception/world-bending demos; screenshots from current deck.")
    footer(s, 7)

    # 10. Theory page, preserve Gaussian illustration and formula images if possible but make layout sane.
    s = prs.slides[9]
    remove_nonpictures(s)
    pics = picture_shapes(s)
    # Keep first picture large enough if it is the Gaussian illustration.
    if pics:
        set_pic(pics[0], 0.95, 1.56, w=3.65)
    title(s, "Splats Are Harder Than Points Under Nonlinear Projection",
          "A Gaussian splat has a screen-space footprint, so changing projection affects more than its center.",
          "Background")
    label(s, 5.0, 1.65, 1.45, 0.52, "mesh\nvertices", BLUE)
    label(s, 6.85, 1.65, 1.45, 0.52, "points\ncenters", RED)
    label(s, 8.7, 1.65, 1.85, 0.52, "splats\ncenter + ellipse", TEAL)
    bullet_box(s, 5.0, 2.65, 5.75, 1.0, [
        "Meshes/points mainly require projecting vertices or centers.",
        "Splats also require projecting a 3D covariance into a 2D ellipse."
    ], 12)
    # Formula pictures are preserved but unknown order; better add text cue, leave existing formula images if present.
    # Reposition remaining pictures to formula row.
    if len(pics) > 1:
        set_pic(pics[1], 4.95, 4.28, w=4.15)
    if len(pics) > 2:
        set_pic(pics[2], 9.15, 4.28, w=2.65)
    text_box(s, 5.0, 5.18, 6.8, 0.34, "J is the projection Jacobian. Under fisheye, J must become the fisheye Jacobian.", 10.2, False, GREY, PP_ALIGN.CENTER)
    source(s, "Formula: standard 3DGS covariance projection; see Kerbl et al., 2023. Visual Gaussian illustration from current deck.")
    footer(s, 10)

    # 11. Architecture diagram.
    s = prs.slides[10]
    clear_slide(s)
    title(s, "Project Architecture",
          "The prototype connects asset conversion, renderer modifications, and desktop/VR preview scenes.",
          "Our Project")
    # pipeline
    for i, (txt, col) in enumerate([("PLY", BLUE), ("SPZ", BLUE), ("SOG", PURPLE)]):
        label(s, 0.85, 1.58 + i * 0.48, 0.95, 0.32, txt, col, size=9.2)
    text_box(s, 0.72, 1.22, 1.4, 0.2, "input formats", 9, True, BLACK, PP_ALIGN.CENTER)
    arrow(s, 1.92, 2.05, 2.45, 2.05)
    label(s, 2.58, 1.65, 1.55, 0.8, "importer /\nconverter", TEAL, size=9.5)
    arrow(s, 4.22, 2.05, 4.72, 2.05)
    label(s, 4.85, 1.65, 1.75, 0.8, "Gaussian\nSplatAsset", TEAL, size=9.5)
    arrow(s, 6.7, 2.05, 7.2, 2.05)
    label(s, 7.33, 1.42, 2.05, 0.55, "UnityGaussianSplatting", GREEN, size=9)
    label(s, 7.33, 2.18, 2.05, 0.55, "FOV / fisheye params", GOLD, size=9)
    arrow(s, 9.48, 2.05, 10.0, 2.05)
    label(s, 10.15, 1.45, 1.65, 0.5, "desktop\npreview", BLUE, size=9)
    label(s, 10.15, 2.2, 1.65, 0.5, "VR / XR\npreview", PURPLE, size=9)
    # methods
    text_box(s, 0.86, 4.0, 2.2, 0.24, "implementation paths", 11.5, True)
    label(s, 0.86, 4.45, 2.25, 0.62, "center-only direct\nfirst attempt", RED, size=8.8)
    label(s, 3.45, 4.45, 2.25, 0.62, "point debug\ncenter validation", GOLD, size=8.8)
    label(s, 6.05, 4.45, 2.25, 0.62, "cubemap fisheye\nbaseline", BLUE, size=8.8)
    label(s, 8.65, 4.45, 2.75, 0.62, "direct covariance-aware\nfisheye", TEAL, size=8.8)
    text_box(s, 1.0, 5.72, 10.8, 0.34, "The presentation follows this implementation story: what broke, what we used to debug it, and how the final methods compare.", 11.5, False, GREY, PP_ALIGN.CENTER)
    footer(s, 11)

    prs.save(OUT)
    strip_effects(OUT)
    print(OUT)


if __name__ == "__main__":
    polish()
