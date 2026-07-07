from pathlib import Path
from zipfile import ZipFile, ZIP_DEFLATED
from tempfile import NamedTemporaryFile
from lxml import etree
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor


ROOT = Path(r"E:\3dgs fisheye projection")
SRC = ROOT / "Presentation" / "FOV_Fisheye_3DGS_clean_LaTeX_formulas_aspect.pptx"
OUT = ROOT / "Presentation" / "FOV_Fisheye_3DGS_polished_opening_architecture.pptx"
ASSETS = ROOT / "Presentation" / "extract_target_images"
FORMULAS = ROOT / "Presentation" / "formula_renders_clean_aspect"


BLACK = RGBColor(18, 18, 18)
GREY = RGBColor(86, 86, 86)
LIGHT_GREY = RGBColor(247, 248, 250)
MID = RGBColor(210, 218, 228)
TEAL = RGBColor(0, 135, 145)
BLUE = RGBColor(43, 93, 173)
RED = RGBColor(196, 65, 62)
PURPLE = RGBColor(113, 82, 165)
GOLD = RGBColor(186, 142, 45)
GREEN = RGBColor(64, 145, 98)


def clear_slide(slide):
    spTree = slide.shapes._spTree
    for shp in list(slide.shapes):
        spTree.remove(shp._element)


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None, section=None):
    if section:
        add_text(slide, 0.72, 0.28, 2.5, 0.22, section.upper(), 8, True, TEAL)
    add_text(slide, 0.72, 0.55, 11.8, 0.46, title, 25, True, BLACK)
    if subtitle:
        add_text(slide, 0.74, 1.02, 11.1, 0.28, subtitle, 10.5, False, GREY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.34), Inches(0.7), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def add_footer(slide, num):
    add_text(slide, 0.72, 7.08, 6.4, 0.18, "3DGS FOV/Fisheye Rendering for VR Scene Exploration", 6.5, False, RGBColor(130, 130, 130))
    add_text(slide, 12.25, 7.08, 0.35, 0.18, str(num), 7, False, RGBColor(130, 130, 130), PP_ALIGN.RIGHT)


def bullet_list(slide, x, y, w, h, items, size=13, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.04)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return box


def label(slide, x, y, w, h, text, line=TEAL, fill=RGBColor(255, 255, 255), size=11, bold=True):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(1.1)
    tf = sh.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = BLACK
    return sh


def arrow(slide, x1, y1, x2, y2, color=GREY, width=1.3):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    c.line.end_arrowhead = True
    return c


def fit_size(image_path, max_w, max_h=None):
    with Image.open(image_path) as im:
        px_w, px_h = im.size
    aspect = px_h / px_w
    fit_w = max_w
    fit_h = fit_w * aspect
    if max_h is not None and fit_h > max_h:
        fit_h = max_h
        fit_w = fit_h / aspect
    return fit_w, fit_h


def add_image(slide, path, x, y, w, h=None):
    fit_w, fit_h = fit_size(path, w, h)
    xoff = (w - fit_w) / 2 if w > fit_w else 0
    pic = slide.shapes.add_picture(str(path), Inches(x + xoff), Inches(y), width=Inches(fit_w), height=Inches(fit_h))
    return pic


def add_source(slide, text):
    add_text(slide, 0.72, 6.78, 11.7, 0.18, text, 6.5, False, RGBColor(125, 125, 125))


def add_formula(slide, name, x, y, w, h=None, caption=None):
    path = FORMULAS / f"{name}.png"
    pic = add_image(slide, path, x, y, w, h)
    if caption:
        add_text(slide, x, y + (h if h else 0.42) + 0.04, w, 0.16, caption, 7.2, False, GREY, PP_ALIGN.CENTER)
    return pic


def strip_visual_effects_from_pptx(path: Path):
    effect_tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)
    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in effect_tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(path)


def build():
    prs = Presentation(SRC)

    # Slide 1
    s = prs.slides[0]
    clear_slide(s)
    add_text(s, 0.72, 0.75, 10.5, 1.0, "FOV and Fisheye Distortion Rendering\nof 3D Gaussian Splats for VR Scene Exploration", 31, True, BLACK)
    add_text(s, 0.74, 2.0, 8.6, 0.32, "A Unity prototype for nonlinear view transformations of Gaussian splat scenes", 14.5, False, GREY)
    for i, (txt, col) in enumerate([("PLY / SPZ / SOG", BLUE), ("Unity + URP", TEAL), ("Desktop + VR", GREEN), ("Cubemap vs Direct", PURPLE)]):
        label(s, 0.74 + i * 2.15, 2.62, 1.72, 0.34, txt, col, RGBColor(255, 255, 255), 9.5, True)
    add_text(s, 0.74, 5.65, 5.6, 0.32, "Yingfangzhong SUN  ·  Supervisor: Daniel Filonik", 12, False, BLACK)
    # Use existing QR if available.
    qr = ASSETS / "slide1_pic9.png"
    if qr.exists():
        add_image(s, qr, 10.25, 4.9, 1.25, 1.25)
        add_text(s, 9.88, 6.22, 2.0, 0.18, "GitHub repository", 8, False, GREY, PP_ALIGN.CENTER)
    # Decorative but content-related tiny pipeline.
    label(s, 7.6, 3.45, 1.15, 0.52, "assets", BLUE)
    arrow(s, 8.78, 3.71, 9.28, 3.71)
    label(s, 9.35, 3.45, 1.15, 0.52, "renderer", TEAL)
    arrow(s, 10.53, 3.71, 11.03, 3.71)
    label(s, 11.1, 3.45, 1.15, 0.52, "VR view", GREEN)

    # Slide 2
    s = prs.slides[1]
    clear_slide(s)
    add_title(s, "Talk Roadmap", "A technical story: from 3DGS representation to nonlinear rendering in Unity/VR.")
    toc = [
        ("1", "Background", "3D Gaussian Splatting and why splats are different from meshes/points."),
        ("2", "Motivation", "Why FOV and fisheye transformations are useful for VR scene exploration."),
        ("3", "Asset Pipeline", "PLY, SPZ, and SOG formats entering a Unity GaussianSplatAsset workflow."),
        ("4", "Implementation Story", "Center-only fisheye, point debugging, cubemap baseline, and direct covariance-aware rendering."),
        ("5", "Comparison", "FOV vs fisheye, cubemap vs direct, desktop vs VR."),
        ("6", "Discussion", "What is implemented, what remains for ISMAR Demo, and related systems.")
    ]
    for i, (num, title, desc) in enumerate(toc):
        y = 1.65 + i * 0.78
        col = [BLUE, TEAL, PURPLE, GOLD, GREEN, RED][i]
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(y + 0.08), Inches(0.06), Inches(0.45))
        bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
        add_text(s, 1.12, y, 0.35, 0.25, num, 15, True, col, PP_ALIGN.RIGHT)
        add_text(s, 1.68, y, 2.35, 0.28, title, 14.5, True, BLACK)
        add_text(s, 4.2, y + 0.02, 7.7, 0.32, desc, 10, False, GREY)
    add_footer(s, 2)

    # Slide 3
    s = prs.slides[2]
    clear_slide(s)
    add_title(s, "What Is 3D Gaussian Splatting?", "3DGS represents a scene as many anisotropic transparent Gaussians and renders them in real time.", "Background")
    bullet_list(s, 0.9, 1.72, 4.15, 2.7, [
        "Input: multi-view images with calibrated camera poses.",
        "Representation: millions of 3D Gaussians with position, opacity, color/SH, scale, and rotation.",
        "Rendering: project Gaussians to screen-space ellipses, sort, and alpha-blend.",
        "Key distinction: a splat has extent, not just a point center."
    ], 13)
    # Pipeline
    labels = [("images\n+ poses", BLUE), ("3D\nGaussians", TEAL), ("2D\nellipses", PURPLE), ("novel\nview", GREEN)]
    x = 5.65
    for i, (txt, col) in enumerate(labels):
        label(s, x + i * 1.55, 2.0, 1.1, 0.62, txt, col, RGBColor(255, 255, 255), 9.5)
        if i < len(labels) - 1:
            arrow(s, x + i * 1.55 + 1.14, 2.31, x + (i + 1) * 1.55 - 0.08, 2.31)
    add_text(s, 5.62, 3.2, 5.25, 0.75, "For nonlinear projection, the crucial object is the projected ellipse footprint.", 15, True, TEAL, PP_ALIGN.CENTER)
    add_source(s, "References: Kerbl et al., 3D Gaussian Splatting for Real-Time Radiance Field Rendering, SIGGRAPH 2023; Mildenhall et al., NeRF, ECCV 2020.")
    add_footer(s, 3)

    # Slide 4
    s = prs.slides[3]
    clear_slide(s)
    add_title(s, "From Captured Images to a Renderable Scene", "3DGS turns posed photos into an explicit set of splats that can be rasterized from new viewpoints.", "Background")
    img = ASSETS / "slide4_pic2.gif"
    if img.exists():
        add_image(s, img, 0.95, 1.55, 5.55, 3.15)
    bullet_list(s, 6.95, 1.65, 4.9, 2.25, [
        "Training reconstructs scene radiance into a cloud of 3D Gaussians.",
        "At runtime, rendering is explicit and fast compared with neural ray marching.",
        "This makes 3DGS attractive for interactive desktop and VR exploration."
    ], 14)
    add_text(s, 7.05, 4.25, 4.8, 0.58, "But: changing the view model is harder than changing a mesh camera.", 16, True, TEAL)
    add_source(s, "Image source: LearnOpenCV 3D Gaussian Splatting explanation; method reference: Kerbl et al., 2023.")
    add_footer(s, 4)

    # Slide 5
    s = prs.slides[4]
    clear_slide(s)
    add_title(s, "Why 3DGS Is Interesting for VR", "Photorealistic real-world scenes can be inspected interactively, but the viewer still has a limited field of view.", "Background")
    img = ASSETS / "slide5_pic4.jpg"
    if img.exists():
        add_image(s, img, 0.82, 1.48, 6.5, 3.65)
    bullet_list(s, 7.65, 1.62, 4.35, 2.4, [
        "Real-time novel-view rendering.",
        "Scene detail comes from captured images rather than manual modeling.",
        "Useful for immersive inspection of real spaces.",
        "Open question: how should nonlinear FOV/fisheye views work for splats?"
    ], 14)
    add_source(s, "Image source: 3DGS web viewer screenshot used in current deck. Replace/verify with original URL if needed.")
    add_footer(s, 5)

    # Slide 6
    s = prs.slides[5]
    clear_slide(s)
    add_title(s, "Motivation: Bring More Scene Content Into View", "FOV and fisheye controls can help reveal peripheral or hidden content during VR scene exploration.", "Motivation")
    bullet_list(s, 0.9, 1.72, 5.1, 2.2, [
        "VR users only see a limited instantaneous field of view.",
        "Important scene content may be peripheral, occluded, or spatially hard to inspect.",
        "FOV and fisheye transformations can bring surrounding context into the current view."
    ], 14)
    label(s, 6.65, 1.75, 2.1, 0.72, "normal view", BLUE)
    arrow(s, 8.9, 2.1, 9.5, 2.1, TEAL, 1.8)
    label(s, 9.68, 1.75, 2.1, 0.72, "wide / fisheye view", TEAL)
    add_text(s, 6.6, 3.25, 5.35, 0.65, "Project goal: interactive nonlinear viewing transformations for Gaussian splat scenes in Unity and VR.", 16, True, TEAL, PP_ALIGN.CENTER)
    add_footer(s, 6)

    # Slide 7
    s = prs.slides[6]
    clear_slide(s)
    add_title(s, "Inspiration: Flexible Spatial Transformations", "We first studied simpler rendering demos, then asked how similar transformations behave for Gaussian splats.", "Motivation")
    img1 = ASSETS / "slide7_pic6.jpg"
    img2 = ASSETS / "slide7_pic7.jpg"
    if img1.exists():
        add_image(s, img1, 0.8, 1.5, 5.25, 2.95)
    if img2.exists():
        add_image(s, img2, 6.85, 1.5, 5.25, 2.95)
    add_text(s, 0.95, 4.65, 4.95, 0.26, "World bending / spatial deformation demo", 11, True, BLACK, PP_ALIGN.CENTER)
    add_text(s, 7.0, 4.65, 4.95, 0.26, "FOV / fisheye projection demo", 11, True, BLACK, PP_ALIGN.CENTER)
    add_text(s, 1.15, 5.45, 10.6, 0.45, "Challenge for 3DGS: splats are not ordinary geometry, so the projection must transform centers and screen-space footprints consistently.", 14, True, TEAL, PP_ALIGN.CENTER)
    add_source(s, "Demo sources: advisor-provided fisheye and inception/world-bending demos; screenshots from current deck.")
    add_footer(s, 7)

    # Slide 10
    s = prs.slides[9]
    clear_slide(s)
    add_title(s, "Splats Are Harder Than Points Under Nonlinear Projection", "A Gaussian splat has a screen-space footprint, so changing projection affects more than its center.", "Background")
    # use original Gaussian image if available
    img = ASSETS / "slide10_pic18.png"
    if img.exists():
        add_image(s, img, 0.95, 1.55, 3.5, 2.55)
    label(s, 4.95, 1.65, 1.55, 0.58, "mesh\nvertices", BLUE)
    label(s, 6.95, 1.65, 1.55, 0.58, "points\ncenters", RED)
    label(s, 8.95, 1.65, 1.75, 0.58, "splats\ncenter + ellipse", TEAL)
    add_text(s, 4.95, 2.65, 5.9, 0.8, "For meshes or point clouds, nonlinear projection mainly moves vertices or centers. For splats, the projected 2D ellipse also changes.", 13.5, False, BLACK)
    add_formula(s, "covariance", 4.95, 4.15, 4.7, 0.45, "Projected covariance")
    add_text(s, 9.9, 4.03, 2.25, 0.7, "J is the projection Jacobian.\nFor fisheye, J must become J_fish.", 10.5, False, GREY)
    add_source(s, "Formula: standard 3DGS covariance projection; see Kerbl et al., 2023. Visual Gaussian illustration from current deck.")
    add_footer(s, 10)

    # Slide 11
    s = prs.slides[10]
    clear_slide(s)
    add_title(s, "Project Architecture", "The prototype connects asset conversion, renderer modifications, and desktop/VR preview scenes.", "Our Project")
    # Input column
    label(s, 0.8, 1.65, 1.1, 0.42, "PLY", BLUE)
    label(s, 0.8, 2.2, 1.1, 0.42, "SPZ", BLUE)
    label(s, 0.8, 2.75, 1.1, 0.42, "SOG", PURPLE)
    add_text(s, 0.62, 1.25, 1.6, 0.2, "Input formats", 10, True, BLACK, PP_ALIGN.CENTER)
    arrow(s, 2.0, 2.42, 2.55, 2.42)
    label(s, 2.65, 1.7, 1.65, 1.15, "Importer /\nconverter", TEAL)
    arrow(s, 4.38, 2.28, 4.9, 2.28)
    label(s, 5.02, 1.7, 1.85, 1.15, "Gaussian\nSplatAsset", TEAL)
    arrow(s, 6.95, 2.28, 7.47, 2.28)
    label(s, 7.6, 1.45, 2.05, 0.62, "UnityGaussianSplatting\nrenderer", GREEN, size=9.2)
    label(s, 7.6, 2.35, 2.05, 0.62, "FOV / fisheye\nparameters", GOLD, size=9.2)
    arrow(s, 9.75, 2.28, 10.25, 2.28)
    label(s, 10.38, 1.42, 1.8, 0.55, "Desktop\npreview", BLUE)
    label(s, 10.38, 2.22, 1.8, 0.55, "VR / XR\npreview", PURPLE)
    # Methods row
    add_text(s, 0.82, 4.05, 2.3, 0.24, "Rendering approaches", 12, True, BLACK)
    label(s, 0.85, 4.5, 2.3, 0.72, "Cubemap fisheye\n6 faces / 12 in VR", BLUE, size=9.5)
    label(s, 3.65, 4.5, 2.55, 0.72, "Direct fisheye\ncenter + covariance", TEAL, size=9.5)
    label(s, 6.75, 4.5, 2.35, 0.72, "Debug point view\ncenter validation", GOLD, size=9.5)
    label(s, 9.65, 4.5, 2.35, 0.72, "Comparison\nquality + cost", GREEN, size=9.5)
    add_text(s, 0.92, 5.75, 11.1, 0.42, "Story of the implementation: start with a simple direct projection, debug center transforms, build a cubemap baseline, then move to covariance-aware direct rendering.", 12.5, False, GREY, PP_ALIGN.CENTER)
    add_footer(s, 11)

    prs.save(OUT)
    strip_visual_effects_from_pptx(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
