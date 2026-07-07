from pathlib import Path
from tempfile import NamedTemporaryFile
from zipfile import ZIP_DEFLATED, ZipFile

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"E:\3dgs fisheye projection")
SRC = ROOT / "Presentation" / "FOV_Fisheye_3DGS_slide12_two_small_tables.pptx"
OUT = ROOT / "Presentation" / "FOV_Fisheye_3DGS_slide12_academic_tables.pptx"

BLACK = RGBColor(20, 24, 30)
GREY = RGBColor(82, 86, 94)
MUTED = RGBColor(118, 124, 132)
TEAL = RGBColor(0, 132, 143)
RULE = RGBColor(42, 47, 55)
MID_RULE = RGBColor(132, 140, 150)
LIGHT_RULE = RGBColor(222, 226, 231)


def clear_slide(slide):
    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        tree.remove(shape._element)


def fit_frame(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.015)
    tf.margin_right = Inches(0.015)
    tf.margin_top = Inches(0.005)
    tf.margin_bottom = Inches(0.005)


def text_box(slide, x, y, w, h, text, size=7, bold=False, color=BLACK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    fit_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def rule(slide, x, y, w, color=RULE, h=0.006):
    return rect(slide, x, y, w, h, color)


def academic_table(slide, x, y, w, caption, columns, rows, widths, table_no):
    text_box(
        slide,
        x,
        y,
        w,
        0.16,
        f"Table {table_no}. {caption}",
        5.7,
        False,
        BLACK,
        PP_ALIGN.CENTER,
        italic=False,
    )

    y0 = y + 0.28
    header_y = y0 + 0.10
    row_h = 0.43

    rule(slide, x, y0, w, RULE, 0.007)
    cx = x
    for label, cw in zip(columns, widths):
        text_box(slide, cx + 0.02, header_y, cw - 0.04, 0.13, label, 5.45, True, BLACK, PP_ALIGN.LEFT)
        cx += cw
    rule(slide, x, y0 + 0.34, w, MID_RULE, 0.005)

    start_y = y0 + 0.44
    for ri, row in enumerate(rows):
        ry = start_y + ri * row_h
        cx = x
        for ci, (value, cw) in enumerate(zip(row, widths)):
            text_box(slide, cx + 0.02, ry + 0.06, cw - 0.04, row_h - 0.08, value, 5.35, ci == 0, BLACK)
            cx += cw
        if ri < len(rows) - 1:
            rule(slide, x, ry + row_h - 0.01, w, LIGHT_RULE, 0.003)
    rule(slide, x, start_y + len(rows) * row_h - 0.01, w, RULE, 0.007)


def footer(slide):
    text_box(slide, 0.54, 5.31, 4.35, 0.14, "3DGS FOV/Fisheye Rendering for VR Scene Exploration", 4.8, False, MUTED)
    text_box(slide, 9.34, 5.31, 0.2, 0.14, "12", 5.0, False, MUTED, PP_ALIGN.RIGHT)


def strip_effects(path):
    tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)
    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(path)


def main():
    prs = Presentation(SRC)
    slide = prs.slides[11]
    clear_slide(slide)

    text_box(slide, 0.54, 0.38, 8.9, 0.36, "Gaussian Splat Formats and Tool Context", 18.5, True)
    text_box(
        slide,
        0.56,
        0.78,
        8.65,
        0.25,
        "We normalize multiple splat representations and situate the prototype relative to existing 3DGS toolchains.",
        8.0,
        False,
        GREY,
    )
    rect(slide, 0.54, 1.10, 0.58, 0.03, TEAL)

    academic_table(
        slide,
        0.72,
        1.46,
        4.00,
        "Supported splat representations.",
        ["Fmt.", "Origin / role", "Prototype handling"],
        [
            ["PLY", "Canonical 3DGS export", "Importer; large-scene preview"],
            ["SPZ", "Compact capture format", "UnityGaussianSplatting path"],
            ["SOG", "PlayCanvas web container", "Importer; color/memory fixes"],
            ["LOD", "Derived preview levels", "Reduced splat budgets for VR"],
        ],
        [0.52, 1.48, 2.00],
        1,
    )

    academic_table(
        slide,
        5.25,
        1.46,
        4.00,
        "Related 3DGS toolchains.",
        ["Tool", "Relevant capability", "Role in this work"],
        [
            ["UnityGS", "Unity renderer/assets", "Base renderer"],
            ["PlayCanvas", "FOV/fisheye splat rendering", "Reference method"],
            ["Nerfstudio", "Training/export pipeline", "Pipeline context"],
        ],
        [0.88, 1.78, 1.34],
        2,
    )

    text_box(
        slide,
        1.05,
        4.74,
        7.45,
        0.20,
        "Together, these tables define the data and software boundary of the demo before introducing the rendering methods.",
        6.6,
        True,
        TEAL,
        PP_ALIGN.CENTER,
    )
    footer(slide)

    prs.save(OUT)
    strip_effects(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
