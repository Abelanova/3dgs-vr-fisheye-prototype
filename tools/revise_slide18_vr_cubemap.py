import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from io import BytesIO
from pathlib import Path
from tempfile import NamedTemporaryFile
from zipfile import ZIP_DEFLATED, ZipFile
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"E:\3dgs fisheye projection")
SRC = ROOT / "Presentation" / "_tmp_slide18_source_copy.pptx"
OUT = ROOT / "Presentation" / "FOV_Fisheye_3DGS_slide18_vr_cubemap_refined.pptx"
FORMULA_DIR = ROOT / "Presentation" / "formula_renders_vr_cubemap"

BLACK = RGBColor(20, 24, 30)
GREY = RGBColor(82, 86, 94)
MUTED = RGBColor(118, 124, 132)
TEAL = RGBColor(0, 132, 143)
PURPLE = RGBColor(113, 82, 165)
RULE = RGBColor(78, 86, 96)
LIGHT = RGBColor(222, 226, 231)
PALE = RGBColor(247, 249, 251)
WHITE = RGBColor(255, 255, 255)


def clear_slide(slide):
    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        tree.remove(shape._element)


def fit_frame(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)


def text_box(slide, x, y, w, h, text, size=8, bold=False, color=BLACK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    fit_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def bullet_box(slide, x, y, w, h, items, size=6.5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    fit_frame(tf)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2.5)
        r = p.add_run()
        r.text = item
        r.font.name = "Aptos"
        r.font.size = Pt(size)
        r.font.color.rgb = BLACK
    return box


def footer(slide):
    text_box(slide, 0.54, 5.31, 4.35, 0.14, "3DGS FOV/Fisheye Rendering for VR Scene Exploration", 4.8, False, MUTED)
    text_box(slide, 9.34, 5.31, 0.2, 0.14, "18", 5.0, False, MUTED, PP_ALIGN.RIGHT)


def rounded_box(slide, x, y, w, h, text, size=6.4, bold=True, line=LIGHT, fill=WHITE, color=BLACK):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(0.8)
    tf = sh.text_frame
    tf.clear()
    fit_frame(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return sh


def arrow(slide, x1, y1, x2, y2, color=RULE):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(0.9)
    conn.line.end_arrowhead = True
    return conn


def rule(slide, x, y, w, color=LIGHT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.005))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def render_formula(filename, formula, fontsize=18):
    FORMULA_DIR.mkdir(parents=True, exist_ok=True)
    path = FORMULA_DIR / filename
    plt.rcParams.update({
        "mathtext.fontset": "stix",
        "font.family": "STIXGeneral",
    })
    fig = plt.figure(figsize=(5.6, 0.55), dpi=220)
    fig.patch.set_alpha(0)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")
    ax.text(0.5, 0.5, formula, ha="center", va="center", fontsize=fontsize, color="#14181e")
    fig.savefig(path, transparent=True, bbox_inches="tight", pad_inches=0.02)
    plt.close(fig)
    return path


def fit_picture(slide, path_or_blob, x, y, max_w, max_h):
    if isinstance(path_or_blob, (bytes, bytearray)):
        source = BytesIO(path_or_blob)
    else:
        source = str(path_or_blob)
    pic = slide.shapes.add_picture(source, Inches(x), Inches(y))
    aspect = pic.width / pic.height if pic.height else 1
    w = max_w
    h = w / aspect
    if h > max_h:
        h = max_h
        w = h * aspect
    pic.width = Inches(w)
    pic.height = Inches(h)
    return pic


def strip_effects(path):
    tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)
    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(path)


def main():
    formula_count = render_formula("vr_capture_counts.png", r"$N_{\mathrm{faces}} = 6\,N_{\mathrm{eyes}}, \qquad N_{\mathrm{RT}} = 2\,N_{\mathrm{faces}}$", 18)
    formula_sample = render_formula("vr_sampling.png", r"$I_e(u,v)=C_e\!\left(\mathrm{face}(\mathbf{d}_e),\,\mathrm{uv}(\mathbf{d}_e)\right)$", 18)

    prs = Presentation(SRC)
    slide = prs.slides[17]  # VR Cubemap Implementation in this deck
    clear_slide(slide)

    text_box(slide, 0.54, 0.38, 8.9, 0.36, "VR Cubemap Fisheye Implementation", 18.5, True)
    text_box(
        slide,
        0.56,
        0.78,
        8.55,
        0.25,
        "The desktop cubemap baseline was extended to XR by maintaining eye-specific capture buffers and compositing the result in a stereo-aware URP pass.",
        7.9,
        False,
        GREY,
    )
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.54), Inches(1.10), Inches(0.58), Inches(0.03))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    # Method diagram.
    text_box(slide, 0.70, 1.34, 3.9, 0.16, "Stereo capture schedule", 7.6, True, TEAL)
    rounded_box(slide, 0.74, 1.68, 0.86, 0.30, "Left eye", 6.2, True, TEAL)
    rounded_box(slide, 0.74, 2.22, 0.86, 0.30, "Right eye", 6.2, True, PURPLE)
    for row_y, line_color in [(1.68, TEAL), (2.22, PURPLE)]:
        for i in range(6):
            rounded_box(slide, 1.92 + i * 0.43, row_y, 0.27, 0.30, f"F{i+1}", 4.7, True, line=line_color)
    arrow(slide, 4.58, 1.95, 5.04, 1.95)
    rounded_box(slide, 5.18, 1.60, 1.18, 0.70, "staging /\nactive buffers", 6.1)
    arrow(slide, 6.50, 1.95, 6.96, 1.95)
    rounded_box(slide, 7.10, 1.60, 1.22, 0.70, "stereo-aware\nURP composite", 5.9)

    # Video placeholder.
    ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.62), Inches(1.40), Inches(0.78), Inches(1.10))
    ph.fill.solid()
    ph.fill.fore_color.rgb = PALE
    ph.line.color.rgb = LIGHT
    ph.line.width = Pt(0.8)
    text_box(slide, 8.68, 1.72, 0.66, 0.42, "VR demo\nvideo", 6.4, True, MUTED, PP_ALIGN.CENTER)

    # Formula and implementation notes.
    fit_picture(slide, formula_count, 0.88, 2.78, 3.35, 0.34)
    text_box(slide, 0.88, 3.22, 3.35, 0.15, "For stereo VR, this becomes 12 cubemap faces and 24 render textures with double buffering.", 5.4, False, GREY, PP_ALIGN.CENTER)
    fit_picture(slide, formula_sample, 0.88, 3.55, 3.35, 0.34)
    text_box(slide, 0.88, 3.98, 3.35, 0.15, "The composite shader samples the cubemap for the current stereo eye.", 5.4, False, GREY, PP_ALIGN.CENTER)

    text_box(slide, 4.75, 2.82, 4.10, 0.18, "Engineering details from the VR branch", 7.8, True)
    bullet_box(slide, 4.75, 3.13, 4.05, 1.05, [
        "Capture work is distributed by face pairs per frame to reduce frame spikes.",
        "Active/staging buffers are swapped only after a complete update cycle.",
        "The output camera is isolated while the URP pass displays the composite.",
        "Mono composite and stereo-separation controls were added for VR diagnostics."
    ], 6.2)

    rule(slide, 0.82, 4.42, 8.16)
    text_box(
        slide,
        0.90,
        4.68,
        8.0,
        0.18,
        "Role in the presentation: this method is the robust VR baseline, but it motivates the later direct renderer because the capture cost scales with eye count and cubemap resolution.",
        6.3,
        True,
        TEAL,
        PP_ALIGN.CENTER,
    )
    footer(slide)

    prs.save(OUT)
    strip_effects(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
