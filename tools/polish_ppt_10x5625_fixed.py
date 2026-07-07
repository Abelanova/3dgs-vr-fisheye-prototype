from pathlib import Path
from tempfile import NamedTemporaryFile
from zipfile import ZIP_DEFLATED, ZipFile

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"E:\3dgs fisheye projection")
SRC = ROOT / "Presentation" / "FOV_Fisheye_3DGS_clean_LaTeX_formulas_aspect.pptx"
OUT = ROOT / "Presentation" / "FOV_Fisheye_3DGS_conservative_10x5625_fixed.pptx"

BLACK = RGBColor(18, 18, 18)
GREY = RGBColor(82, 82, 82)
MUTED = RGBColor(120, 120, 120)
TEAL = RGBColor(0, 135, 145)
BLUE = RGBColor(43, 93, 173)
RED = RGBColor(196, 65, 62)
PURPLE = RGBColor(113, 82, 165)
GREEN = RGBColor(64, 145, 98)
GOLD = RGBColor(186, 142, 45)
WHITE = RGBColor(255, 255, 255)


def remove_nonpictures(slide):
    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        if not hasattr(shape, "image"):
            tree.remove(shape._element)


def clear_slide(slide):
    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        tree.remove(shape._element)


def fit_frame(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)


def text_box(slide, x, y, w, h, text, size=10, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    fit_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def bullet_box(slide, x, y, w, h, items, size=9.0, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    fit_frame(tf)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = item
        r.font.name = "Aptos"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def title(slide, heading, subtitle=None, section=None):
    if section:
        text_box(slide, 0.50, 0.18, 1.6, 0.16, section.upper(), 6.8, True, TEAL)
    text_box(slide, 0.50, 0.42, 8.9, 0.38, heading, 18.5, True)
    if subtitle:
        text_box(slide, 0.51, 0.85, 8.7, 0.24, subtitle, 8.2, False, GREY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.50), Inches(1.18), Inches(0.55), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def footer(slide, n):
    text_box(slide, 0.50, 5.40, 4.4, 0.12, "3DGS FOV/Fisheye Rendering for VR Scene Exploration", 4.8, False, MUTED)
    text_box(slide, 9.25, 5.40, 0.25, 0.12, str(n), 5.2, False, MUTED, PP_ALIGN.RIGHT)


def source(slide, text):
    text_box(slide, 0.50, 5.18, 8.9, 0.16, text, 4.8, False, MUTED)


def label(slide, x, y, w, h, text, line=TEAL, size=7.2):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = line
    sh.line.width = Pt(0.9)
    tf = sh.text_frame
    tf.clear()
    fit_frame(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = BLACK
    return sh


def arrow(slide, x1, y1, x2, y2, color=GREY):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(1)
    conn.line.end_arrowhead = True
    return conn


def pictures(slide):
    return [s for s in slide.shapes if hasattr(s, "image")]


def fit_pic(shape, x, y, max_w, max_h):
    aspect = shape.width / shape.height if shape.height else 1
    w = max_w
    h = w / aspect
    if h > max_h:
        h = max_h
        w = h * aspect
    shape.left = Inches(x)
    shape.top = Inches(y)
    shape.width = Inches(w)
    shape.height = Inches(h)


def strip_effects(path):
    tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)
    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(path)


def polish():
    prs = Presentation(SRC)

    # 1. Title slide
    s = prs.slides[0]
    remove_nonpictures(s)
    title_slots = [(0.75, 4.42, 0.75, 0.42), (1.85, 4.42, 0.75, 0.42), (3.15, 4.18, 0.78, 0.78), (4.30, 4.18, 0.78, 0.78), (8.15, 4.10, 0.72, 0.72)]
    for pic, slot in zip(pictures(s), title_slots):
        fit_pic(pic, *slot)
    text_box(s, 0.58, 0.58, 8.8, 0.82, "FOV and Fisheye Distortion Rendering\nof 3D Gaussian Splats for VR Scene Exploration", 22, True)
    text_box(s, 0.60, 1.64, 8.3, 0.22, "A Unity prototype for nonlinear view transformations of Gaussian splat scenes", 10.5, False, GREY)
    for i, (txt, col) in enumerate([("PLY / SPZ / SOG", BLUE), ("Unity + URP", TEAL), ("Desktop + VR", GREEN), ("Cubemap vs Direct", PURPLE)]):
        label(s, 0.60 + i * 2.1, 2.28, 1.55, 0.28, txt, col, 7.4)
    text_box(s, 0.60, 5.08, 5.1, 0.18, "Yingfangzhong SUN  /  Supervisor: Daniel Filonik", 8.0, False)

    # 2. Roadmap
    s = prs.slides[1]
    clear_slide(s)
    title(s, "Talk Roadmap", "From 3DGS representation to nonlinear rendering in Unity/VR.")
    rows = [
        ("1", "Background", "3DGS representation and splat footprints."),
        ("2", "Motivation", "Why FOV/fisheye helps VR scene exploration."),
        ("3", "Formats", "PLY, SPZ, and SOG in a Unity asset pipeline."),
        ("4", "Implementation", "Direct attempt, point debug, cubemap, direct covariance."),
        ("5", "Comparison", "FOV vs fisheye; cubemap vs direct; desktop vs VR."),
        ("6", "Discussion", "Current status, next steps, and references."),
    ]
    for i, (num, head, desc) in enumerate(rows):
        y = 1.45 + i * 0.52
        text_box(s, 0.78, y, 0.26, 0.18, num, 11.5, True, [BLUE, TEAL, PURPLE, GOLD, GREEN, RED][i], PP_ALIGN.RIGHT)
        text_box(s, 1.30, y, 1.55, 0.18, head, 10.5, True)
        text_box(s, 3.30, y + 0.01, 5.7, 0.18, desc, 8.2, False, GREY)
    footer(s, 2)

    # 3. Keep existing diagram, clean only explanatory text.
    s = prs.slides[2]
    for sh in list(s.shapes):
        if hasattr(sh, "text") and sh.text.strip() and ("Core idea" in sh.text or "3DGS is" in sh.text or "input:" in sh.text or "References:" in sh.text):
            sh.text_frame.clear()
        if hasattr(sh, "text") and sh.top > Inches(4.95):
            s.shapes._spTree.remove(sh._element)
    text_box(s, 0.62, 1.25, 2.1, 0.18, "Core idea", 10.5, True, TEAL)
    bullet_box(s, 0.62, 1.62, 2.3, 1.12, ["Input: posed images.", "Representation: anisotropic 3D Gaussians.", "Rendering: project, sort, alpha-blend."], 8.2)
    text_box(s, 0.62, 3.04, 2.35, 0.34, "A splat has position, opacity, color, and shape.", 7.8, True, TEAL, PP_ALIGN.CENTER)
    source(s, "References: Kerbl et al., 3D Gaussian Splatting, SIGGRAPH 2023; Mildenhall et al., NeRF, ECCV 2020.")

    # 4. 3DGS intro
    s = prs.slides[3]
    remove_nonpictures(s)
    pics = pictures(s)
    if pics:
        fit_pic(pics[0], 0.65, 1.45, 5.25, 3.35)
    title(s, "From Captured Images to a Renderable Scene", "3DGS learns an explicit splat representation that can be rendered from new viewpoints.", "Background")
    bullet_box(s, 6.15, 1.62, 3.15, 1.25, ["Learns radiance from posed images.", "Renders explicit splats interactively.", "Good fit for desktop and VR inspection."], 8.8)
    text_box(s, 6.15, 3.55, 3.15, 0.46, "But the camera model also changes the splat footprint.", 10.2, True, TEAL)
    source(s, "Image source: LearnOpenCV 3D Gaussian Splatting explanation. Method reference: Kerbl et al., 2023.")
    footer(s, 4)

    # 5. VR motivation
    s = prs.slides[4]
    remove_nonpictures(s)
    pics = pictures(s)
    if pics:
        fit_pic(pics[0], 0.65, 1.45, 5.35, 3.45)
    title(s, "Why 3DGS Is Interesting for VR", "Photorealistic captured spaces can be inspected interactively, but the user still has limited instantaneous view.", "Background")
    bullet_box(s, 6.28, 1.62, 3.1, 1.7, ["Real-time novel views.", "Detail comes from captured spaces.", "Useful for immersive inspection.", "Question: how should wide FOV and fisheye work for splats?"], 8.5)
    source(s, "Image source: 3DGS web viewer screenshot from current deck. Verify original URL before final submission.")
    footer(s, 5)

    # 6. Motivation
    s = prs.slides[5]
    remove_nonpictures(s)
    title(s, "Motivation: Bring More Scene Content Into View", "FOV and fisheye controls can reveal peripheral or hidden content during VR scene exploration.", "Motivation")
    bullet_box(s, 0.68, 1.55, 3.65, 1.4, ["VR has limited instantaneous field of view.", "Relevant content may lie outside the current view.", "Nonlinear viewing can compress more context into the display."], 9.2)
    label(s, 5.05, 1.62, 1.35, 0.44, "normal view", BLUE, 7.8)
    arrow(s, 6.52, 1.84, 7.00, 1.84, TEAL)
    label(s, 7.18, 1.62, 1.7, 0.44, "wide / fisheye view", TEAL, 7.8)
    text_box(s, 4.8, 3.05, 4.35, 0.54, "Goal: interactive nonlinear viewing transformations for Gaussian splat scenes in Unity and VR.", 10.4, True, TEAL, PP_ALIGN.CENTER)
    footer(s, 6)

    # 7. Inspirations
    s = prs.slides[6]
    remove_nonpictures(s)
    pics = pictures(s)
    if len(pics) >= 2:
        fit_pic(pics[0], 0.65, 1.42, 4.1, 2.35)
        fit_pic(pics[1], 5.2, 1.42, 4.1, 2.35)
    title(s, "Inspiration: Flexible Spatial Transformations", "We first studied simpler demos, then asked how similar transformations behave for Gaussian splats.", "Motivation")
    text_box(s, 0.8, 3.9, 3.8, 0.2, "World bending / spatial deformation demo", 8.2, True, BLACK, PP_ALIGN.CENTER)
    text_box(s, 5.35, 3.9, 3.8, 0.2, "FOV / fisheye projection demo", 8.2, True, BLACK, PP_ALIGN.CENTER)
    text_box(s, 0.9, 4.48, 8.1, 0.36, "Challenge for 3DGS: centers and screen-space footprints must transform consistently.", 9.8, True, TEAL, PP_ALIGN.CENTER)
    source(s, "Demo sources: advisor-provided fisheye and inception/world-bending demos; screenshots from current deck.")
    footer(s, 7)

    # 10. Theory
    s = prs.slides[9]
    remove_nonpictures(s)
    pics = pictures(s)
    if pics:
        fit_pic(pics[0], 0.65, 1.36, 3.15, 2.25)
    title(s, "Splats Are Harder Than Points Under Nonlinear Projection", "A Gaussian splat has a screen-space footprint, so changing projection affects more than its center.", "Background")
    label(s, 4.15, 1.45, 1.1, 0.42, "mesh\nvertices", BLUE, 7.2)
    label(s, 5.55, 1.45, 1.1, 0.42, "points\ncenters", RED, 7.2)
    label(s, 6.95, 1.45, 1.55, 0.42, "splats\ncenter + ellipse", TEAL, 7.0)
    bullet_box(s, 4.15, 2.20, 4.75, 0.75, ["Meshes/points: project vertices or centers.", "Splats: also project 3D covariance to a 2D ellipse."], 8.4)
    if len(pics) > 1:
        fit_pic(pics[1], 4.15, 3.35, 2.95, 0.55)
    if len(pics) > 2:
        fit_pic(pics[2], 7.38, 3.35, 1.75, 0.55)
    text_box(s, 4.15, 4.18, 4.75, 0.28, "J is the projection Jacobian. Under fisheye, J becomes the fisheye Jacobian.", 7.6, False, GREY, PP_ALIGN.CENTER)
    source(s, "Formula: standard 3DGS covariance projection; see Kerbl et al., 2023. Visual Gaussian illustration from current deck.")
    footer(s, 10)

    # 11. Architecture
    s = prs.slides[10]
    clear_slide(s)
    title(s, "Project Architecture", "The prototype connects asset conversion, renderer modifications, and desktop/VR preview scenes.", "Our Project")
    for i, (txt, col) in enumerate([("PLY", BLUE), ("SPZ", BLUE), ("SOG", PURPLE)]):
        label(s, 0.62, 1.48 + i * 0.32, 0.62, 0.22, txt, col, 6.4)
    text_box(s, 0.48, 1.18, 1.0, 0.15, "input formats", 6.8, True, BLACK, PP_ALIGN.CENTER)
    arrow(s, 1.38, 1.82, 1.74, 1.82)
    label(s, 1.88, 1.48, 1.15, 0.6, "importer /\nconverter", TEAL, 7.0)
    arrow(s, 3.12, 1.82, 3.50, 1.82)
    label(s, 3.64, 1.48, 1.18, 0.6, "Gaussian\nSplatAsset", TEAL, 7.0)
    arrow(s, 4.92, 1.82, 5.30, 1.82)
    label(s, 5.43, 1.34, 1.55, 0.42, "UnityGaussianSplatting", GREEN, 6.5)
    label(s, 5.43, 1.95, 1.55, 0.42, "FOV / fisheye params", GOLD, 6.5)
    arrow(s, 7.10, 1.82, 7.48, 1.82)
    label(s, 7.60, 1.34, 1.15, 0.38, "desktop\npreview", BLUE, 6.5)
    label(s, 7.60, 1.98, 1.15, 0.38, "VR / XR\npreview", PURPLE, 6.5)
    text_box(s, 0.62, 3.28, 1.65, 0.18, "implementation paths", 8.8, True)
    label(s, 0.62, 3.70, 1.75, 0.44, "center-only direct\nfirst attempt", RED, 6.4)
    label(s, 2.62, 3.70, 1.75, 0.44, "point debug\ncenter validation", GOLD, 6.4)
    label(s, 4.62, 3.70, 1.75, 0.44, "cubemap fisheye\nbaseline", BLUE, 6.4)
    label(s, 6.62, 3.70, 2.0, 0.44, "direct covariance-aware\nfisheye", TEAL, 6.2)
    text_box(s, 0.82, 4.72, 7.95, 0.28, "The presentation follows this story: what broke, how we debugged it, and how the final methods compare.", 8.3, False, GREY, PP_ALIGN.CENTER)
    footer(s, 11)

    prs.save(OUT)
    strip_effects(OUT)
    print(OUT)


if __name__ == "__main__":
    polish()
