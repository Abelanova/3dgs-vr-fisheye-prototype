from io import BytesIO
from pathlib import Path
from tempfile import NamedTemporaryFile
from zipfile import ZIP_DEFLATED, ZipFile

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"E:\3dgs fisheye projection")
SRC = ROOT / "Presentation" / "FOV_Fisheye_3DGS_clean_LaTeX_formulas_aspect.pptx"
OUT = ROOT / "Presentation" / "FOV_Fisheye_3DGS_slide16_cubemap_refined.pptx"

BLACK = RGBColor(20, 24, 30)
GREY = RGBColor(82, 86, 94)
MUTED = RGBColor(118, 124, 132)
TEAL = RGBColor(0, 132, 143)
RULE = RGBColor(90, 98, 108)
LIGHT = RGBColor(222, 226, 231)
PALE = RGBColor(247, 249, 251)
WHITE = RGBColor(255, 255, 255)


def clear_slide(slide):
    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        tree.remove(shape._element)


def fit_frame(tf):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)


def text_box(slide, x, y, w, h, text, size=8, bold=False, color=BLACK, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    fit_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def footer(slide):
    text_box(slide, 0.54, 5.31, 4.35, 0.14, "3DGS FOV/Fisheye Rendering for VR Scene Exploration", 4.8, False, MUTED)
    text_box(slide, 9.34, 5.31, 0.2, 0.14, "16", 5.0, False, MUTED, PP_ALIGN.RIGHT)


def rounded_box(slide, x, y, w, h, text, size=7.0, bold=True, fill=WHITE):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = LIGHT
    sh.line.width = Pt(0.8)
    tf = sh.text_frame
    tf.clear()
    fit_frame(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = BLACK
    return sh


def arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = RULE
    conn.line.width = Pt(0.9)
    conn.line.end_arrowhead = True
    return conn


def rule(slide, x, y, w, color=LIGHT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.005))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def picture_blobs(slide):
    return [shape.image.blob for shape in slide.shapes if hasattr(shape, "image")]


def add_fit_picture(slide, blob, x, y, max_w, max_h):
    pic = slide.shapes.add_picture(BytesIO(blob), Inches(x), Inches(y))
    aspect = pic.width / pic.height if pic.height else 1
    w = max_w
    h = w / aspect
    if h > max_h:
        h = max_h
        w = h * aspect
    pic.width = Inches(w)
    pic.height = Inches(h)
    return pic


def strip_effects(path):
    tags = {
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectLst",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}effectDag",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}outerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}innerShdw",
        "{http://schemas.openxmlformats.org/drawingml/2006/main}prstShdw",
    }
    with NamedTemporaryFile(delete=False, suffix=".pptx", dir=path.parent) as tmp:
        tmp_path = Path(tmp.name)
    with ZipFile(path, "r") as zin, ZipFile(tmp_path, "w", ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.endswith(".xml") and (
                item.filename.startswith("ppt/slides/")
                or item.filename.startswith("ppt/slideMasters/")
                or item.filename.startswith("ppt/slideLayouts/")
                or item.filename.startswith("ppt/theme/")
            ):
                try:
                    root = etree.fromstring(data)
                    changed = False
                    for node in list(root.iter()):
                        if node.tag in tags:
                            parent = node.getparent()
                            if parent is not None:
                                parent.remove(node)
                                changed = True
                    if changed:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)
    tmp_path.replace(path)


def main():
    prs = Presentation(SRC)
    slide = prs.slides[15]
    blobs = picture_blobs(slide)
    # Original order: formula 1, formula 2, cubemap diagram, paper figure.
    formula1, formula2, cube, paper_fig = blobs[0], blobs[1], blobs[2], blobs[3]

    clear_slide(slide)
    text_box(slide, 0.54, 0.38, 8.9, 0.36, "Attempt 2: Cubemap-Based Fisheye Rendering", 18.5, True)
    text_box(
        slide,
        0.56,
        0.78,
        8.45,
        0.25,
        "Instead of projecting splats directly through a nonlinear camera, render perspective cubemap faces and perform fisheye sampling in a composite pass.",
        7.9,
        False,
        GREY,
    )
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.54), Inches(1.10), Inches(0.58), Inches(0.03))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    # Top method pipeline.
    for i, face in enumerate(["+X", "-X", "+Y", "-Y", "+Z", "-Z"]):
        rounded_box(slide, 0.68 + i * 0.55, 1.48, 0.42, 0.28, face, 5.6)
    text_box(slide, 0.72, 1.88, 3.10, 0.14, "six 90° perspective captures", 5.4, False, GREY, PP_ALIGN.CENTER)
    arrow(slide, 4.02, 1.63, 4.56, 1.63)
    rounded_box(slide, 4.72, 1.38, 1.36, 0.50, "cubemap\ntextures", 6.4)
    arrow(slide, 6.22, 1.63, 6.76, 1.63)
    rounded_box(slide, 6.92, 1.38, 1.28, 0.50, "fisheye\ncomposite", 6.4)
    add_fit_picture(slide, cube, 8.36, 1.16, 1.10, 1.10)

    # Explanation blocks.
    text_box(slide, 0.82, 2.45, 3.75, 0.18, "Rendering stage", 8.0, True)
    text_box(
        slide,
        0.82,
        2.76,
        3.75,
        0.72,
        "Each capture uses the ordinary perspective 3DGS renderer. This avoids changing the splat footprint model during the first pass.",
        7.1,
        False,
        BLACK,
    )
    text_box(slide, 5.05, 2.45, 3.75, 0.18, "Composite stage", 8.0, True)
    text_box(
        slide,
        5.05,
        2.76,
        3.75,
        0.72,
        "For each output pixel, invert the fisheye projection to a ray direction, then sample the corresponding cubemap face.",
        7.1,
        False,
        BLACK,
    )

    # Formula and compact visual reference.
    rule(slide, 0.82, 3.72, 8.05)
    add_fit_picture(slide, formula1, 1.10, 3.92, 2.55, 0.42)
    add_fit_picture(slide, formula2, 1.12, 4.42, 2.55, 0.42)
    text_box(slide, 3.85, 4.03, 2.35, 0.18, "screen pixel → fisheye ray → cubemap color", 6.8, False, BLACK, PP_ALIGN.CENTER)
    add_fit_picture(slide, paper_fig, 6.48, 3.82, 2.10, 0.92)
    text_box(slide, 6.48, 4.78, 2.10, 0.12, "ray-based fisheye/cubemap mapping reference", 4.6, False, MUTED, PP_ALIGN.CENTER)

    text_box(
        slide,
        0.90,
        5.03,
        8.10,
        0.18,
        "Role in our project: a stable baseline for comparison; the cost increases in VR because stereo output requires separate cubemap captures.",
        6.5,
        True,
        TEAL,
        PP_ALIGN.CENTER,
    )
    footer(slide)

    prs.save(OUT)
    strip_effects(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
